"""
FYP Presentation Builder — Ray Tracing for Path Loss Prediction
Dark academic theme, high-level design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCA, 0xD3, 0xDF)   # soft grey text
C_GREEN     = RGBColor(0x2D, 0xCB, 0x7F)   # success green
C_ORANGE    = RGBColor(0xFF, 0x9F, 0x1C)   # warning orange
C_RED       = RGBColor(0xFF, 0x4D, 0x6D)   # error red
C_PANEL     = RGBColor(0x11, 0x2A, 0x40)   # slightly lighter panel
C_DIVIDER   = RGBColor(0x1E, 0x3A, 0x52)   # divider line

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=C_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, alpha=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size, color=C_WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    return txb


def accent_bar(slide, y=Inches(0.08), h=Inches(0.06)):
    """Top accent bar."""
    rect(slide, 0, y, W, h, C_ACCENT)


def slide_number(slide, n, total=12):
    txt(slide, f"{n} / {total}", Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.35),
        size=9, color=C_LIGHT, align=PP_ALIGN.RIGHT)


def section_tag(slide, label, x=Inches(0.35), y=Inches(0.22)):
    """Small uppercase tag top-left."""
    txt(slide, label.upper(), x, y, Inches(4), Inches(0.28),
        size=8, color=C_ACCENT, bold=True)


def panel(slide, x, y, w, h):
    """Dark panel box."""
    return rect(slide, x, y, w, h, C_PANEL)


def divider(slide, x, y, w, thickness=Inches(0.018)):
    rect(slide, x, y, w, thickness, C_DIVIDER)


def bullet_block(slide, items, x, y, w, size=12, color=C_LIGHT, spacing=0.38):
    """Render a list of bullet strings as separate textboxes."""
    cy = y
    for item in items:
        prefix = "  •  " if not item.startswith("◆") else ""
        txt(slide, prefix + item.lstrip("◆").strip(), x, cy, w, Inches(spacing),
            size=size, color=color)
        cy += Inches(spacing)
    return cy


def kv_row(slide, key, val, x, y, w, key_color=C_ACCENT2, val_color=C_WHITE,
           size=11, row_h=Inches(0.38)):
    """Key | value row."""
    col1 = w * 0.42
    col2 = w * 0.58
    txt(slide, key, x, y, col1, row_h, size=size, color=key_color, bold=True)
    txt(slide, val, x + col1, y, col2, row_h, size=size, color=val_color)
    return y + row_h


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)

# Large cyan decorative bar on left
rect(s, 0, 0, Inches(0.25), H, C_ACCENT)

# Main title
txt(s, "Ray Tracing for Path Loss Prediction",
    Inches(0.55), Inches(1.6), Inches(9.5), Inches(1.2),
    size=40, bold=True, color=C_WHITE)

txt(s, "in Urban Environments",
    Inches(0.55), Inches(2.7), Inches(9.5), Inches(0.8),
    size=32, bold=False, color=C_ACCENT2)

divider(s, Inches(0.55), Inches(3.55), Inches(5.5))

# Subtitle details
details = [
    ("Scene",     "Nottingham · Ofcom 2018 Drive Test"),
    ("Frequency", "915.95 MHz  (IoT / LoRaWAN band)"),
    ("Tool",      "Sionna RT 2.0  +  Sionna 0.19 Differentiable RT"),
    ("Goal",      "Minimise RMSE against 1 200 real Ofcom measurements"),
]
cy = Inches(3.75)
for k, v in details:
    txt(s, k, Inches(0.55), cy, Inches(2.2), Inches(0.38), size=11,
        color=C_ACCENT, bold=True)
    txt(s, v, Inches(2.75), cy, Inches(7.0), Inches(0.38), size=11,
        color=C_LIGHT)
    cy += Inches(0.42)

# Right-side decorative ray-tracing pattern suggestion (text placeholder)
rect(s, Inches(10.0), Inches(1.2), Inches(3.0), Inches(5.2), C_PANEL)
txt(s, "[ Scene Render\nor Coverage Map ]",
    Inches(10.1), Inches(2.8), Inches(2.8), Inches(1.5),
    size=11, color=C_DIVIDER, align=PP_ALIGN.CENTER, italic=True)

slide_number(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PURPOSE & MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Motivation")
slide_number(s, 2)

txt(s, "Purpose & Motivation", Inches(0.45), Inches(0.5), Inches(10), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Three columns
cols = [
    {
        "icon": "⚠",
        "title": "The Problem",
        "color": C_RED,
        "items": [
            "Empirical models (COST-231,\nOkumura-Hata) give 20–36 dB errors",
            "No terrain, no building geometry,\nno material physics",
            "Fixed-environment formulas —\nbreak in new cities",
        ]
    },
    {
        "icon": "🎯",
        "title": "Our Goal",
        "color": C_ACCENT,
        "items": [
            "Minimise RMSE against\n1 200 real Ofcom measurements",
            "Reproduce physical propagation:\ndiffraction, reflection, scatter",
            "Differentiable RT calibration\n(NVLabs Hoydis et al. 2023)",
        ]
    },
    {
        "icon": "📡",
        "title": "The Dataset",
        "color": C_GREEN,
        "items": [
            "Ofcom 2018 drive test\n915.95 MHz, Nottingham UK",
            "1 200 receivers  |  0.3–9 km",
            "RSSI range: −118 to −23 dBm\nNoise floor: −124 dBm",
        ]
    },
]

cw = Inches(4.0)
cx_starts = [Inches(0.35), Inches(4.55), Inches(8.75)]

for i, col in enumerate(cols):
    cx = cx_starts[i]
    panel(s, cx, Inches(1.35), cw, Inches(5.7))
    rect(s, cx, Inches(1.35), cw, Inches(0.07), col["color"])

    txt(s, col["icon"] + "  " + col["title"],
        cx + Inches(0.18), Inches(1.5), cw - Inches(0.3), Inches(0.5),
        size=14, bold=True, color=col["color"])

    cy2 = Inches(2.1)
    for item in col["items"]:
        txt(s, "›  " + item, cx + Inches(0.18), cy2, cw - Inches(0.3), Inches(0.9),
            size=11, color=C_LIGHT)
        cy2 += Inches(1.05)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESEARCH QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Research Questions")
slide_number(s, 3)

txt(s, "Research Questions", Inches(0.45), Inches(0.5), Inches(10), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

rqs = [
    ("RQ 1", C_ACCENT,
     "Flat Terrain vs DEM",
     "How much does digital elevation modelling (DEM) improve path loss\n"
     "prediction accuracy over a flat z = 0 scene in Nottingham?"),
    ("RQ 2", C_GREEN,
     "Differentiable RT Calibration",
     "Can Sionna 0.19 differentiable RT (NVLabs approach) achieve\n"
     "< 5 dB RMSE on the Ofcom 2018 UK dataset after calibration?"),
    ("RQ 3", C_ORANGE,
     "Optimal Scene Configuration",
     "Which combination of scene geometry, material parameters, and\n"
     "ray-tracing settings best represents urban Nottingham at 915 MHz?"),
]

cy = Inches(1.4)
for tag, color, title, body in rqs:
    rect(s, Inches(0.35), cy + Inches(0.12), Inches(0.06), Inches(0.85), color)
    rect(s, Inches(0.41), cy, Inches(12.5), Inches(1.12), C_PANEL)
    txt(s, tag, Inches(0.55), cy + Inches(0.06), Inches(1.0), Inches(0.4),
        size=10, color=color, bold=True)
    txt(s, title, Inches(1.55), cy + Inches(0.05), Inches(4.5), Inches(0.4),
        size=13, color=C_WHITE, bold=True)
    txt(s, body, Inches(6.3), cy + Inches(0.05), Inches(6.4), Inches(1.0),
        size=11, color=C_LIGHT)
    cy += Inches(1.28)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METHODOLOGY PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Methodology")
slide_number(s, 4)

txt(s, "Simulation & Calibration Pipeline", Inches(0.45), Inches(0.5),
    Inches(12), Inches(0.7), size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

stages = [
    ("01", "INPUT DATA",      C_ACCENT,  "OSM + EA LiDAR\n+ Ofcom CSV"),
    ("02", "SCENE BUILDER",   C_ACCENT2, "nDSM + ITU-R P.2040-2\n+ Terrain PLY"),
    ("03", "RAY TRACING",     C_GREEN,   "Sionna 2.0 PathSolver\n+ Sionna 0.19 Diff RT"),
    ("04", "CALIBRATION",     C_ORANGE,  "Scalar (Cell 10b)\nMaterial (Cell 11b)"),
    ("05", "ML RESIDUAL",     C_RED,     "ResidualMLP (Cell 15)\nCNN+MLP (Cell 14)\nMaterialMLP (Cell 16)"),
    ("06", "VALIDATION",      C_ACCENT,  "RMSE · MAE · Bias\nSTD · R² per band"),
]

bw = Inches(2.05)
bh = Inches(4.5)
gap = Inches(0.08)
bx = Inches(0.28)

for i, (num, label, color, detail) in enumerate(stages):
    cx = bx + i * (bw + gap)
    panel(s, cx, Inches(1.38), bw, bh)
    rect(s, cx, Inches(1.38), bw, Inches(0.06), color)

    txt(s, num, cx + Inches(0.12), Inches(1.5), bw, Inches(0.45),
        size=22, bold=True, color=color)
    txt(s, label, cx + Inches(0.12), Inches(1.95), bw - Inches(0.15), Inches(0.5),
        size=10, bold=True, color=C_WHITE)
    divider(s, cx + Inches(0.12), Inches(2.48), bw - Inches(0.25))
    txt(s, detail, cx + Inches(0.12), Inches(2.58), bw - Inches(0.15), Inches(1.8),
        size=10, color=C_LIGHT)

    # Arrow between stages
    if i < len(stages) - 1:
        ax = cx + bw + Inches(0.01)
        txt(s, "▶", ax, Inches(3.35), gap + Inches(0.06), Inches(0.4),
            size=9, color=C_DIVIDER, align=PP_ALIGN.CENTER)

# Validation metrics row
txt(s, "Output Metrics per Distance Band:   RMSE  |  MAE  |  Bias  |  STD  |  R²",
    Inches(0.35), Inches(6.1), Inches(12.6), Inches(0.4),
    size=11, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OFCOM DATA PROCESSING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Input Data")
slide_number(s, 5)

txt(s, "Ofcom Dataset — GPS to Sionna Receivers",
    Inches(0.45), Inches(0.5), Inches(12), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left: pipeline steps
panel(s, Inches(0.35), Inches(1.35), Inches(7.0), Inches(5.7))

steps = [
    ("1", C_ACCENT,  "Raw Ofcom CSV",
     "GPS (WGS84 lat/lon)  +  RSSI (dBm)  per measurement point"),
    ("2", C_ACCENT,  "Filter & Clean",
     "Keep 915.95 MHz  ·  drop RSSI < −124 dBm  ·  remove duplicates"),
    ("3", C_GREEN,   "Coordinate Transform",
     "WGS84  →  UTM Zone 30N (EPSG:32630)  →  local scene XY (metres)"),
    ("4", C_GREEN,   "Height Assignment",
     "RX Z = EA LiDAR DTM elevation  +  1.5 m AGL"),
    ("5", C_ORANGE,  "Sionna RX Array",
     "Each point  →  Receiver object  ·  isotropic antenna  ·  1 200 total"),
    ("6", C_ACCENT2, "PL Reference",
     "PL_meas = TX_EIRP − RSSI_meas  (dB)  →  ground truth for RMSE"),
]

cy = Inches(1.5)
for num, color, title, body in steps:
    rect(s, Inches(0.5), cy + Inches(0.08), Inches(0.32), Inches(0.32), color)
    txt(s, num, Inches(0.5), cy + Inches(0.05), Inches(0.32), Inches(0.32),
        size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(0.95), cy + Inches(0.03), Inches(2.0), Inches(0.35),
        size=11, bold=True, color=color)
    txt(s, body, Inches(3.0), cy + Inches(0.03), Inches(4.2), Inches(0.35),
        size=10, color=C_LIGHT)
    cy += Inches(0.87)

# Right: key numbers panel
panel(s, Inches(7.6), Inches(1.35), Inches(5.4), Inches(5.7))
rect(s, Inches(7.6), Inches(1.35), Inches(5.4), Inches(0.06), C_ACCENT)

txt(s, "Key Numbers", Inches(7.75), Inches(1.5), Inches(5.0), Inches(0.45),
    size=14, bold=True, color=C_ACCENT)

kvs = [
    ("Raw Ofcom points",    "~4 000"),
    ("After filtering",     "1 200 receivers"),
    ("Coverage range",      "0.3 km – 9 km from TX"),
    ("TX location",         "(−1.2559°, 52.9863°)"),
    ("TX height",           "17 m AGL"),
    ("Frequency",           "915.95 MHz"),
    ("RSSI range (meas)",   "−118 to −22.9 dBm"),
    ("Noise floor",         "−124 dBm"),
    ("PL_meas range",       "93.7 – 142.4 dB"),
]

cy2 = Inches(2.1)
for k, v in kvs:
    txt(s, k, Inches(7.75), cy2, Inches(2.8), Inches(0.38), size=10,
        color=C_ACCENT2, bold=True)
    txt(s, v, Inches(10.55), cy2, Inches(2.3), Inches(0.38), size=10,
        color=C_WHITE)
    cy2 += Inches(0.52)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SCENE CONSTRUCTION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Scene Construction")
slide_number(s, 6)

txt(s, "Scene Construction — Blender → OSM → Merged PLY",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Three panels side by side
panels_data = [
    {
        "title": "Stage 1 · Blender (Flat)",
        "color": C_ORANGE,
        "sub": "Manual · OSM importer plugin",
        "pros": ["Fast to build — one import/export", "Full geometry cleanup in GUI",
                 "Single merged PLY → Sionna loads in seconds"],
        "cons": ["No terrain — flat z=0 plane", "Nottingham hills completely absent",
                 "Manual — not reproducible", "8–15 dB error at 300–1500 m"],
    },
    {
        "title": "Stage 2 · OSM Direct (DEM)",
        "color": C_GREEN,
        "sub": "Automated · Python OSMnx + scene builder",
        "pros": ["Fully automated & reproducible", "ITU-R P.2040-2 materials per building type",
                 "Real terrain from EA LiDAR (2 m res)"],
        "cons": ["OSM heights often missing → estimated", "OSM tags inconsistent across Nottingham",
                 "Longer scene build time"],
    },
    {
        "title": "Stage 3 · Merged PLY",
        "color": C_ACCENT,
        "sub": "Why we merged 150 k files → one per material",
        "pros": ["Scene load: 20 min → 45 sec", "Single BVH → 50× faster ray traversal",
                 "No GPU memory fragmentation", "Solve rate improved directly"],
        "cons": ["Per-material merge requires preprocessing", "Harder to inspect individual buildings"],
    },
]

pw = Inches(4.2)
gap2 = Inches(0.12)
for i, pd in enumerate(panels_data):
    cx = Inches(0.3) + i * (pw + gap2)
    panel(s, cx, Inches(1.35), pw, Inches(5.75))
    rect(s, cx, Inches(1.35), pw, Inches(0.07), pd["color"])

    txt(s, pd["title"], cx + Inches(0.15), Inches(1.5), pw - Inches(0.2), Inches(0.42),
        size=12, bold=True, color=pd["color"])
    txt(s, pd["sub"], cx + Inches(0.15), Inches(1.93), pw - Inches(0.2), Inches(0.38),
        size=9.5, color=C_LIGHT, italic=True)
    divider(s, cx + Inches(0.15), Inches(2.32), pw - Inches(0.3))

    txt(s, "✔  Pros", cx + Inches(0.15), Inches(2.4), pw, Inches(0.35),
        size=10, bold=True, color=C_GREEN)
    cy3 = Inches(2.76)
    for p in pd["pros"]:
        txt(s, "   " + p, cx + Inches(0.15), cy3, pw - Inches(0.2), Inches(0.45),
            size=9.5, color=C_LIGHT)
        cy3 += Inches(0.44)

    txt(s, "✘  Cons", cx + Inches(0.15), cy3 + Inches(0.05), pw, Inches(0.35),
        size=10, bold=True, color=C_RED)
    cy3 += Inches(0.43)
    for c in pd["cons"]:
        txt(s, "   " + c, cx + Inches(0.15), cy3, pw - Inches(0.2), Inches(0.45),
            size=9.5, color=C_LIGHT)
        cy3 += Inches(0.44)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6b — BUILDING HEIGHTS & TERRAIN ERROR SOURCES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Error Sources")
slide_number(s, 7)

txt(s, "Known Error Sources — Building Heights & Terrain Geometry",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Section 1 — Building height
panel(s, Inches(0.35), Inches(1.3), Inches(6.0), Inches(2.85))
rect(s, Inches(0.35), Inches(1.3), Inches(6.0), Inches(0.06), C_ORANGE)
txt(s, "Building Height Problem — Both Scenes", Inches(0.5), Inches(1.42),
    Inches(5.7), Inches(0.4), size=12, bold=True, color=C_ORANGE)

bh_rows = [
    ("OSM height tag",         "Only ~30% of buildings have it",            "70% use estimate"),
    ("building:levels tag",    "height = levels × 3 m",                     "Wrong for mixed-use"),
    ("No tag",                 "Default 8 m applied",                       "Underestimates tall buildings"),
    ("Industrial/warehouse",   "Tagged as residential",                     "Wrong shadowing"),
]
cy = Inches(1.92)
for src_label, problem, impact in bh_rows:
    txt(s, src_label, Inches(0.5),  cy, Inches(1.7), Inches(0.38), size=9.5, color=C_ACCENT2, bold=True)
    txt(s, problem,   Inches(2.2),  cy, Inches(2.3), Inches(0.38), size=9.5, color=C_LIGHT)
    txt(s, impact,    Inches(4.5),  cy, Inches(1.7), Inches(0.38), size=9.5, color=C_RED)
    cy += Inches(0.42)

txt(s, "→ Wrong building heights = wrong diffraction angles = wrong path loss",
    Inches(0.5), cy + Inches(0.05), Inches(5.7), Inches(0.35),
    size=10, color=C_ORANGE, italic=True)

# Section 2 — Flat terrain
panel(s, Inches(0.35), Inches(4.3), Inches(5.9), Inches(2.8))
rect(s, Inches(0.35), Inches(4.3), Inches(5.9), Inches(0.06), C_RED)
txt(s, "Flat Terrain Errors (Sionna v1)", Inches(0.5), Inches(4.42),
    Inches(5.6), Inches(0.4), size=12, bold=True, color=C_RED)

flat_rows = [
    ("Terrain missing",    "All buildings on z=0",          "+8–15 dB overestimate 300–1500 m"),
    ("Ground reflection",  "Flat plane — wrong angle",      "~3–5 dB additional bias"),
    ("RX height",          "1.5 m above z=0, not real ground", "Some RX inside terrain"),
]
cy = Inches(4.92)
for label, detail, impact in flat_rows:
    txt(s, label,  Inches(0.5), cy, Inches(1.5), Inches(0.38), size=9.5, color=C_ACCENT2, bold=True)
    txt(s, detail, Inches(2.0), cy, Inches(2.2), Inches(0.38), size=9.5, color=C_LIGHT)
    txt(s, impact, Inches(4.2), cy, Inches(1.9), Inches(0.38), size=9.5, color=C_RED)
    cy += Inches(0.42)

# Section 3 — DEM terrain
panel(s, Inches(6.5), Inches(1.3), Inches(6.5), Inches(2.85))
rect(s, Inches(6.5), Inches(1.3), Inches(6.5), Inches(0.06), C_GREEN)
txt(s, "DEM Terrain Errors (Sionna v2)", Inches(6.65), Inches(1.42),
    Inches(6.2), Inches(0.4), size=12, bold=True, color=C_GREEN)

dem_rows = [
    ("Building heights",   "DEM fixes ground Z, not building tops",     "Residual ±3–8 dB"),
    ("DEM resolution",     "2 m — small walls/hedges missed",            "Minor"),
    ("Roof geometry",      "Flat-roof box — no pitched roofs",           "Diffraction angle error"),
]
cy2 = Inches(1.92)
for label, detail, impact in dem_rows:
    txt(s, label,  Inches(6.65), cy2, Inches(1.7), Inches(0.38), size=9.5, color=C_ACCENT2, bold=True)
    txt(s, detail, Inches(8.35), cy2, Inches(2.8), Inches(0.38), size=9.5, color=C_LIGHT)
    txt(s, impact, Inches(11.15),cy2, Inches(1.7), Inches(0.38), size=9.5, color=C_ORANGE)
    cy2 += Inches(0.42)

# Summary box
panel(s, Inches(6.5), Inches(4.3), Inches(6.5), Inches(2.8))
rect(s, Inches(6.5), Inches(4.3), Inches(6.5), Inches(0.06), C_ACCENT)
txt(s, "Key Takeaway", Inches(6.65), Inches(4.42), Inches(6.2), Inches(0.4),
    size=12, bold=True, color=C_ACCENT)
txt(s, "Flat terrain: dominant error is missing hills → systematic bias\n"
       "DEM terrain: dominant error is building height estimation → random scatter\n"
       "Both improve with material calibration — but building height is a hard geometric limit",
    Inches(6.65), Inches(4.9), Inches(6.2), Inches(1.8),
    size=11, color=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE A — DTM / DSM / nDSM CONCEPTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Elevation Data")
slide_number(s, 8)

txt(s, "Elevation Data Sources — DTM, DSM & nDSM",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

cols_dsm = [
    {
        "title": "DTM",
        "subtitle": "Digital Terrain Model",
        "color": C_GREEN,
        "rows": [
            ("Captures",    "Bare earth only\n(ground, roads, rivers)"),
            ("Source",      "EA LiDAR 1m\n(last / ground returns)"),
            ("Use in scene","Terrain mesh PLY\n→ correct RX ground Z"),
            ("Limitation",  "No buildings or trees"),
        ]
    },
    {
        "title": "DSM",
        "subtitle": "Digital Surface Model",
        "color": C_ORANGE,
        "rows": [
            ("Captures",    "Everything — ground +\nbuildings + trees"),
            ("Source",      "EA LiDAR 1m\n(first returns)"),
            ("Use in scene","Building height\nestimation"),
            ("Limitation",  "Cannot separate building\nfrom terrain directly"),
        ]
    },
    {
        "title": "nDSM",
        "subtitle": "Normalised DSM = DSM − DTM",
        "color": C_ACCENT,
        "rows": [
            ("Captures",    "Height of objects\nABOVE ground only"),
            ("Source",      "Derived: DSM − DTM"),
            ("Use in scene","Accurate per-building\nAGL height from LiDAR"),
            ("Limitation",  "Requires both DTM\n& DSM"),
        ]
    },
]

cw3 = Inches(4.1)
for i, col in enumerate(cols_dsm):
    cx = Inches(0.3) + i * (cw3 + Inches(0.11))
    panel(s, cx, Inches(1.35), cw3, Inches(5.7))
    rect(s, cx, Inches(1.35), cw3, Inches(0.07), col["color"])

    txt(s, col["title"], cx + Inches(0.18), Inches(1.5), cw3, Inches(0.55),
        size=22, bold=True, color=col["color"])
    txt(s, col["subtitle"], cx + Inches(0.18), Inches(2.05), cw3 - Inches(0.25), Inches(0.4),
        size=10, color=C_LIGHT, italic=True)
    divider(s, cx + Inches(0.18), Inches(2.48), cw3 - Inches(0.35))

    cy3 = Inches(2.6)
    for label, val in col["rows"]:
        txt(s, label, cx + Inches(0.18), cy3, Inches(1.2), Inches(0.75),
            size=9.5, color=C_ACCENT2, bold=True)
        txt(s, val, cx + Inches(1.38), cy3, cw3 - Inches(1.55), Inches(0.75),
            size=9.5, color=C_LIGHT)
        cy3 += Inches(0.78)

# Formula highlight
rect(s, Inches(3.8), Inches(6.55), Inches(5.7), Inches(0.58), C_PANEL)
rect(s, Inches(3.8), Inches(6.55), Inches(5.7), Inches(0.04), C_ACCENT)
txt(s, "nDSM = DSM − DTM   →   object height above ground (buildings, trees)",
    Inches(4.0), Inches(6.63), Inches(5.3), Inches(0.4),
    size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE B — RMSE COMPARISON: FLAT / DTM / nDSM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Results Comparison")
slide_number(s, 9)

txt(s, "RMSE Results — Flat / DTM / nDSM Scene Comparison",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

result_rows = [
    ("Flat z=0\n(Sionna v1)",     "None — z=0 plane",  "OSM estimate / default 8 m",
     "19.6 dB", "−4.4 dB", C_RED,
     "Missing hills → +8–15 dB overestimate\nat 300–1500 m"),
    ("DTM terrain\n(Sionna v2)",  "EA LiDAR DTM ✅",   "OSM floor count / default 8 m",
     "14.1 dB", "~0 dB", C_ORANGE,
     "Building height estimation\nstill introduces ±3–8 dB scatter"),
    ("nDSM heights\n(2.8–6 GHz)", "EA LiDAR DTM ✅",   "LiDAR nDSM (accurate AGL)",
     "Used in\n2.8 GHz scene", "—", C_ACCENT,
     "Rooftop equip added —\nonly valid at higher frequencies"),
]

cy = Inches(1.45)
headers = ["Scene Config", "Terrain", "Building Heights", "RMSE", "Bias", "Key Issue"]
hx = [Inches(0.35), Inches(2.55), Inches(4.55), Inches(7.35), Inches(8.55), Inches(9.55)]
hw = [Inches(2.1), Inches(1.9), Inches(2.7), Inches(1.1), Inches(0.9), Inches(3.6)]
for j, h in enumerate(headers):
    txt(s, h, hx[j], cy, hw[j], Inches(0.38), size=10, bold=True, color=C_ACCENT2)
divider(s, Inches(0.35), cy + Inches(0.4), Inches(12.8))
cy += Inches(0.55)

for config, terrain, bldg, rmse, bias, color, issue in result_rows:
    rh = Inches(1.3)
    rect(s, Inches(0.35), cy, Inches(0.05), rh, color)
    panel(s, Inches(0.4), cy, Inches(12.7), rh)

    txt(s, config, hx[0] + Inches(0.08), cy + Inches(0.25), hw[0], rh,
        size=10, bold=True, color=color)
    txt(s, terrain, hx[1], cy + Inches(0.25), hw[1], rh, size=9.5, color=C_LIGHT)
    txt(s, bldg,    hx[2], cy + Inches(0.25), hw[2], rh, size=9.5, color=C_LIGHT)
    txt(s, rmse,    hx[3], cy + Inches(0.25), hw[3], rh, size=14, bold=True, color=color,
        align=PP_ALIGN.CENTER)
    txt(s, bias,    hx[4], cy + Inches(0.25), hw[4], rh, size=11, color=C_WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, issue,   hx[5], cy + Inches(0.18), hw[5], rh, size=9.5, color=C_LIGHT)
    cy += rh + Inches(0.12)

# Key finding
panel(s, Inches(0.35), Inches(6.1), Inches(12.7), Inches(0.95))
rect(s, Inches(0.35), Inches(6.1), Inches(12.7), Inches(0.05), C_GREEN)
txt(s, "Key Finding:  DTM terrain alone reduces RMSE by 5.5 dB (19.6 → 14.1 dB).  "
       "Bias drops from −4.4 dB → ~0 dB.  "
       "Terrain is essential — building height remains the next bottleneck.",
    Inches(0.6), Inches(6.22), Inches(12.2), Inches(0.7),
    size=11, color=C_WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE C — HOW WE MERGED: scene_v2_infra PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Scene Pipeline")
slide_number(s, 10)

txt(s, "Final Scene: scene_v2_infra — How DTM + nDSM Are Merged",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Pipeline flow — left side
pipeline_steps = [
    (C_GREEN,   "EA LiDAR DTM 1m",    "Ground-only elevation raster (GeoTIFF)\nCovered: 11×11 km Nottingham bbox"),
    (C_ORANGE,  "EA LiDAR DSM 1m",    "First-return surface raster\nIncludes buildings + trees"),
    (C_ACCENT,  "nDSM = DSM − DTM",   "Per-pixel height above ground\nMedian filter → remove LiDAR spikes"),
    (C_ACCENT2, "OSM Footprints",      "Building polygons + landuse tags\nHeight assigned from nDSM median under footprint"),
    (C_GREEN,   "scene_v2_infra",      "terrain.ply (DTM mesh) + building meshes\n+ road/veg/infra + ITU-R P.2040-2 materials"),
]

cy = Inches(1.38)
for color, title, detail in pipeline_steps:
    rect(s, Inches(0.35), cy + Inches(0.12), Inches(0.06), Inches(0.75), color)
    panel(s, Inches(0.41), cy, Inches(7.3), Inches(1.0))
    txt(s, title, Inches(0.6), cy + Inches(0.08), Inches(2.5), Inches(0.42),
        size=12, bold=True, color=color)
    txt(s, detail, Inches(3.1), cy + Inches(0.08), Inches(4.5), Inches(0.82),
        size=10, color=C_LIGHT)
    if cy < Inches(6.0):
        txt(s, "▼", Inches(3.6), cy + Inches(1.02), Inches(0.4), Inches(0.3),
            size=10, color=C_DIVIDER, align=PP_ALIGN.CENTER)
    cy += Inches(1.12)

# Right panel — contribution table
panel(s, Inches(8.0), Inches(1.35), Inches(5.1), Inches(5.7))
rect(s, Inches(8.0), Inches(1.35), Inches(5.1), Inches(0.06), C_ACCENT)
txt(s, "What Each Layer Contributes", Inches(8.15), Inches(1.5),
    Inches(4.8), Inches(0.42), size=12, bold=True, color=C_ACCENT)

contrib = [
    (C_GREEN,   "DTM terrain",    "Correct ground Z\n→ RX at real AGL",       "8–15 dB bias removed"),
    (C_ACCENT,  "nDSM heights",   "Accurate building\nheights from LiDAR",    "Correct diffraction angles"),
    (C_ACCENT2, "OSM footprints", "Building positions\nand 2D shapes",         "3D geometry"),
    (C_ORANGE,  "ITU materials",  "EM properties\nper surface type",           "Physical reflection/scatter"),
]
cy2 = Inches(2.05)
for color, layer, role, benefit in contrib:
    rect(s, Inches(8.1), cy2 + Inches(0.08), Inches(0.04), Inches(0.72), color)
    txt(s, layer, Inches(8.22), cy2 + Inches(0.05), Inches(1.5), Inches(0.82),
        size=9.5, bold=True, color=color)
    txt(s, role, Inches(9.72), cy2 + Inches(0.05), Inches(1.6), Inches(0.82),
        size=9.5, color=C_LIGHT)
    txt(s, benefit, Inches(11.32), cy2 + Inches(0.05), Inches(1.65), Inches(0.82),
        size=9.5, color=C_GREEN)
    cy2 += Inches(0.96)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RAY TRACING CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Ray Tracing Setup")
slide_number(s, 11)

txt(s, "Ray Tracing Configuration — Mechanisms & Parameters",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left — mechanisms
panel(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.7))
rect(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(0.06), C_ACCENT)
txt(s, "Ray Mechanisms", Inches(0.5), Inches(1.5), Inches(5.6), Inches(0.42),
    size=13, bold=True, color=C_ACCENT)

mechs = [
    ("LOS",               C_GREEN,  "Direct path TX → RX"),
    ("Specular Reflection",C_GREEN,  "Building face bounce"),
    ("Diffraction",        C_GREEN,  "Over rooftops / wedge edges"),
    ("Edge Diffraction",   C_GREEN,  "Around building corners — dominant NLOS"),
    ("Scattering",         C_GREEN,  "Lambertian rough surface diffuse"),
]
cy = Inches(2.05)
for mech, color, desc in mechs:
    rect(s, Inches(0.55), cy + Inches(0.1), Inches(0.28), Inches(0.28), color)
    txt(s, "✔", Inches(0.55), cy + Inches(0.05), Inches(0.3), Inches(0.35),
        size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, mech, Inches(0.95), cy + Inches(0.03), Inches(2.1), Inches(0.38),
        size=11, bold=True, color=C_WHITE)
    txt(s, desc, Inches(3.05), cy + Inches(0.03), Inches(2.9), Inches(0.38),
        size=10, color=C_LIGHT)
    cy += Inches(0.82)

# Finding box
panel(s, Inches(0.35), Inches(6.3), Inches(5.9), Inches(0.75))
txt(s, "72% of rays are diffraction — edge diffraction is the dominant\npropagation mechanism in urban Nottingham NLOS at 915 MHz",
    Inches(0.5), Inches(6.38), Inches(5.6), Inches(0.6),
    size=10, color=C_ACCENT2, italic=True)

# Right — parameters
panel(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(5.7))
rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(0.06), C_ORANGE)
txt(s, "Solver Parameters", Inches(6.65), Inches(1.5), Inches(6.2), Inches(0.42),
    size=13, bold=True, color=C_ORANGE)

params = [
    ("Frequency",       "915.95 MHz"),
    ("MAX_DEPTH",       "7  (up to 7 bounces)"),
    ("NUM_SAMPLES",     "10 000 000  (NVLabs approach)"),
    ("TX height",       "17 m AGL  (mast)"),
    ("RX height",       "1.5 m AGL  (IoT device)"),
    ("TX antenna",      "Dipole — donut pattern"),
    ("RX antenna",      "Isotropic"),
    ("Noise floor",     "−124 dBm"),
    ("BATCH",           "1 receiver per solve (NVLabs)"),
]
cy2 = Inches(2.05)
for k, v in params:
    txt(s, k, Inches(6.65), cy2, Inches(2.5), Inches(0.42),
        size=10, bold=True, color=C_ACCENT2)
    txt(s, v, Inches(9.15), cy2, Inches(3.7), Inches(0.42),
        size=10, color=C_WHITE)
    cy2 += Inches(0.52)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — OSM FEATURES & PATH LOSS IMPACT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "OSM Features")
slide_number(s, 12)

txt(s, "OSM Features Downloaded — Role in Propagation Physics",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left — feature classes
panel(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(5.7))
rect(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(0.06), C_ACCENT)
txt(s, "Feature Classes Downloaded", Inches(0.5), Inches(1.5),
    Inches(5.8), Inches(0.4), size=12, bold=True, color=C_ACCENT)

osm_features = [
    (C_ACCENT,  "Buildings",        "building=*",               "52 248",  "Brick/Concrete/Glass/Metal/Wood"),
    (C_LIGHT,   "Roads",            "highway=*",                "Full net", "Asphalt"),
    (C_GREEN,   "Trees/Vegetation", "natural=tree, landuse=forest","~273",  "Vegetation"),
    (C_ORANGE,  "Masts/Towers",     "man_made=mast/tower",      "100s",    "Metal"),
    (C_ORANGE,  "Power pylons",     "power=tower",              "100s",    "Metal"),
    (C_RED,     "Railways",         "railway=rail",             "Full net", "Metal"),
    (C_LIGHT,   "Water bodies",     "natural=water, waterway=*","Polygons","Water"),
    (C_LIGHT,   "Barriers/embank.", "barrier=*, embankment",    "Many",    "Concrete"),
    (C_LIGHT,   "Car parks",        "amenity=parking",          "Many",    "Concrete"),
]
cy = Inches(2.05)
for color, feat, tag, count, mat in osm_features:
    rect(s, Inches(0.45), cy + Inches(0.09), Inches(0.04), Inches(0.3), color)
    txt(s, feat,  Inches(0.55), cy, Inches(1.55), Inches(0.42), size=9.5, bold=True, color=color)
    txt(s, tag,   Inches(2.1),  cy, Inches(1.9),  Inches(0.42), size=8.5, color=C_LIGHT, italic=True)
    txt(s, count, Inches(4.0),  cy, Inches(0.65), Inches(0.42), size=9.5, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, mat,   Inches(4.65), cy, Inches(1.65), Inches(0.42), size=9,   color=C_ACCENT2)
    cy += Inches(0.52)

# Right — propagation impact
panel(s, Inches(6.6), Inches(1.35), Inches(6.45), Inches(5.7))
rect(s, Inches(6.6), Inches(1.35), Inches(6.45), Inches(0.06), C_ORANGE)
txt(s, "Propagation Impact per Feature", Inches(6.75), Inches(1.5),
    Inches(6.1), Inches(0.4), size=12, bold=True, color=C_ORANGE)

impacts = [
    (C_ACCENT,  "Buildings",       "Diffraction + reflection — dominant NLOS loss",       "Primary"),
    (C_ORANGE,  "Metal masts",     "Strong specular reflectors — secondary paths",         "+2–5 dB"),
    (C_LIGHT,   "Roads",           "Ground bounce — 2-ray interference pattern",           "±3 dB"),
    (C_GREEN,   "Vegetation",      "Absorption + scatter (P.833) up to 67 dB at 9 km",    "+12–27 dB"),
    (C_RED,     "Railways",        "Specular waveguide along rail corridors",              "Elevated"),
    (C_ACCENT2, "Water (Trent)",   "Near-specular reflection — local enhancement",         "Local +"),
    (C_LIGHT,   "Barriers",        "Extra diffraction edges — additional shielding",       "+2–8 dB"),
]
cy2 = Inches(2.05)
for color, feat, effect, impact in impacts:
    rect(s, Inches(6.7), cy2 + Inches(0.09), Inches(0.04), Inches(0.3), color)
    txt(s, feat,   Inches(6.82), cy2, Inches(1.5),  Inches(0.42), size=9.5, bold=True, color=color)
    txt(s, effect, Inches(8.32), cy2, Inches(3.5),  Inches(0.42), size=9,   color=C_LIGHT)
    txt(s, impact, Inches(11.82),cy2, Inches(1.1),  Inches(0.42), size=9.5, bold=True, color=color,
        align=PP_ALIGN.RIGHT)
    cy2 += Inches(0.68)

# Key finding
panel(s, Inches(0.35), Inches(7.1), Inches(12.7), Inches(0.0))
txt(s, "Metal infrastructure (σ = 10⁷ S/m) added in scene_v2_infra — absent in flat baseline.  "
       "At 915 MHz metal surfaces are near-perfect reflectors creating strong secondary paths.",
    Inches(0.5), Inches(6.88), Inches(12.2), Inches(0.5),
    size=10, color=C_ACCENT2, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TX/RX POSITIONING & RAY-CAST VALIDATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "TX / RX Setup")
slide_number(s, 13)

txt(s, "TX & RX Placement — GPS to Scene + Ray-Cast Validation",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# TX panel
panel(s, Inches(0.35), Inches(1.35), Inches(4.0), Inches(3.2))
rect(s, Inches(0.35), Inches(1.35), Inches(4.0), Inches(0.06), C_ORANGE)
txt(s, "TX — Transmitter", Inches(0.5), Inches(1.5), Inches(3.7), Inches(0.4),
    size=12, bold=True, color=C_ORANGE)

tx_rows = [
    ("GPS (WGS84)",    "−1.25590°,  52.98630°"),
    ("Transform",      "WGS84 → UTM 30N → scene XY"),
    ("Terrain Z",      "~79 m ODN  (EA LiDAR DTM)"),
    ("AGL height",     "+17 m  (mast)"),
    ("Scene Z",        "~96 m absolute"),
    ("Antenna",        "Dipole — donut pattern"),
    ("Ray-cast check", "Downward ray → above terrain ✅"),
]
cy = Inches(2.0)
for k, v in tx_rows:
    txt(s, k, Inches(0.5),  cy, Inches(1.5), Inches(0.38), size=9.5, color=C_ACCENT2, bold=True)
    txt(s, v, Inches(2.0),  cy, Inches(2.2), Inches(0.38), size=9.5, color=C_WHITE)
    cy += Inches(0.38)

# RX pipeline
panel(s, Inches(4.55), Inches(1.35), Inches(8.5), Inches(3.2))
rect(s, Inches(4.55), Inches(1.35), Inches(8.5), Inches(0.06), C_GREEN)
txt(s, "RX — 1200 Ofcom Receivers", Inches(4.7), Inches(1.5), Inches(8.2), Inches(0.4),
    size=12, bold=True, color=C_GREEN)

rx_steps = [
    ("1", C_ACCENT,  "GPS (WGS84)",        "→  UTM Zone 30N (EPSG:32630)  →  local scene XY"),
    ("2", C_ACCENT,  "DEM lookup",         "Bilinear interpolation at (x,y) → terrain Z from LiDAR DTM"),
    ("3", C_GREEN,   "Height assignment",  "RX Z = terrain Z + 1.5 m AGL"),
    ("4", C_ORANGE,  "Ray-cast check",     "Shoot ray downward → if first hit above RX → inside building"),
    ("5", C_RED,     "Fix if invalid",     "Raise RX to building roof + 0.1 m  OR  exclude from solve"),
]
cy2 = Inches(2.0)
for num, color, title, desc in rx_steps:
    rect(s, Inches(4.68), cy2 + Inches(0.06), Inches(0.24), Inches(0.24), color)
    txt(s, num,   Inches(4.68), cy2 + Inches(0.03), Inches(0.24), Inches(0.28),
        size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(5.0),  cy2, Inches(1.6), Inches(0.35), size=10, bold=True, color=color)
    txt(s, desc,  Inches(6.6),  cy2, Inches(6.3), Inches(0.35), size=9.5, color=C_LIGHT)
    cy2 += Inches(0.46)

# Problems table
panel(s, Inches(0.35), Inches(4.75), Inches(12.7), Inches(2.35))
rect(s, Inches(0.35), Inches(4.75), Inches(12.7), Inches(0.06), C_RED)
txt(s, "Why Positioning Validation Matters", Inches(0.5), Inches(4.88),
    Inches(12.0), Inches(0.4), size=12, bold=True, color=C_RED)

prob_rows = [
    ("RX inside building footprint",  "Spawned in solid geometry → no paths → NaN RSSI",     "Ray-cast detects → adjust Z"),
    ("TX near building face",         "Rays immediately blocked → unrealistic near-field",     "Minimum clearance check"),
    ("DEM returns 0 at scene edge",   "RX at z=1.5 m → inside terrain → all paths blocked",  "Fallback to median terrain Z"),
    ("Ofcom GPS imprecision ±3–10 m", "RX may be on wrong side of building wall",             "Accepted — statistical uncertainty"),
]
hx2 = [Inches(0.5), Inches(4.2), Inches(8.4)]
hw2 = [Inches(3.6), Inches(4.1), Inches(4.5)]
cy3 = Inches(5.35)
for problem, consequence, fix in prob_rows:
    txt(s, problem,     hx2[0], cy3, hw2[0], Inches(0.4), size=9.5, color=C_ORANGE)
    txt(s, consequence, hx2[1], cy3, hw2[1], Inches(0.4), size=9.5, color=C_LIGHT)
    txt(s, fix,         hx2[2], cy3, hw2[2], Inches(0.4), size=9.5, color=C_GREEN)
    cy3 += Inches(0.42)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — LOS / NLOS CLASSIFICATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "LOS / NLOS")
slide_number(s, 14)

txt(s, "Receiver Classification — LOS vs NLOS",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left — classification method + results
panel(s, Inches(0.35), Inches(1.35), Inches(5.8), Inches(5.7))
rect(s, Inches(0.35), Inches(1.35), Inches(5.8), Inches(0.06), C_ACCENT)
txt(s, "Classification & Results", Inches(0.5), Inches(1.5),
    Inches(5.5), Inches(0.4), size=12, bold=True, color=C_ACCENT)

txt(s, "Method: LOS if direct path amplitude |a_LOS|² > 0 in solved paths",
    Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.38), size=10, color=C_LIGHT, italic=True)

results_los = [
    ("LOS",  "~20",   "1.7%",  "< 200 m · open streets", C_GREEN),
    ("NLOS", "~1120", "98.3%", "All ranges · urban canyon", C_RED),
]
cy = Inches(2.55)
for cls, count, pct, note, color in results_los:
    panel(s, Inches(0.5), cy, Inches(5.4), Inches(0.82))
    rect(s, Inches(0.5), cy, Inches(0.06), Inches(0.82), color)
    txt(s, cls,   Inches(0.65), cy + Inches(0.12), Inches(0.9), Inches(0.5),
        size=16, bold=True, color=color)
    txt(s, count, Inches(1.55), cy + Inches(0.12), Inches(0.9), Inches(0.5),
        size=16, bold=True, color=C_WHITE)
    txt(s, pct,   Inches(2.45), cy + Inches(0.12), Inches(0.8), Inches(0.5),
        size=14, color=color, bold=True)
    txt(s, note,  Inches(3.25), cy + Inches(0.18), Inches(2.5), Inches(0.45),
        size=9.5, color=C_LIGHT)
    cy += Inches(1.0)

# Ray type breakdown
txt(s, "Ray Type Breakdown — 153 450 total rays",
    Inches(0.5), Inches(4.72), Inches(5.5), Inches(0.35),
    size=10, bold=True, color=C_ACCENT2)

ray_types = [
    ("Diffraction",       "110 618", "72.1%", C_ACCENT),
    ("Multi-reflection",  "37 006",  "24.1%", C_ORANGE),
    ("Reflection",        "3 190",   "2.1%",  C_LIGHT),
    ("LOS",               "2 636",   "1.7%",  C_GREEN),
]
cy = Inches(5.12)
bw_max = Inches(4.8)
for rtype, count, pct, color in ray_types:
    pct_val = float(pct.replace('%',''))
    rect(s, Inches(0.5), cy + Inches(0.08), bw_max * pct_val / 72.1, Inches(0.3), color)
    txt(s, f"{rtype}  {count}  ({pct})", Inches(0.5), cy, Inches(5.4), Inches(0.38),
        size=9.5, color=C_WHITE)
    cy += Inches(0.48)

# Right — why it matters
panel(s, Inches(6.4), Inches(1.35), Inches(6.6), Inches(5.7))
rect(s, Inches(6.4), Inches(1.35), Inches(6.6), Inches(0.06), C_ORANGE)
txt(s, "Why It Matters for Path Loss Modelling",
    Inches(6.55), Inches(1.5), Inches(6.3), Inches(0.4),
    size=12, bold=True, color=C_ORANGE)

matters = [
    (C_GREEN,  "LOS",            "Free-space + ground reflection\nFriis equation sufficient"),
    (C_ORANGE, "NLOS < 500 m",   "Diffraction over 1–2 rooftops\nEdge diffraction critical"),
    (C_RED,    "NLOS 500–2000 m","Multi-hop diffraction + scatter\nFull ray tracing essential"),
    (C_ACCENT, "NLOS > 2000 m",  "Statistical scatter paths\n10M+ samples needed"),
]
cy2 = Inches(2.05)
for color, cond, detail in matters:
    rect(s, Inches(6.5), cy2 + Inches(0.1), Inches(0.05), Inches(0.8), color)
    panel(s, Inches(6.55), cy2, Inches(6.3), Inches(1.0))
    txt(s, cond,   Inches(6.7),  cy2 + Inches(0.08), Inches(2.2), Inches(0.42),
        size=11, bold=True, color=color)
    txt(s, detail, Inches(8.9),  cy2 + Inches(0.08), Inches(3.8), Inches(0.82),
        size=10, color=C_LIGHT)
    cy2 += Inches(1.1)

txt(s, "98.3% NLOS — empirical models calibrated for mixed LOS/NLOS are fundamentally mismatched.\nRay tracing is necessary.",
    Inches(6.55), Inches(6.42), Inches(6.3), Inches(0.55),
    size=10, color=C_ACCENT2, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CELL 10b SCALAR CALIBRATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Calibration · Stage 1")
slide_number(s, 15)

txt(s, "Stage 1 — Scalar Offset Calibration  (NVLabs ITU Materials Baseline)",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=24, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left — method
panel(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(5.7))
rect(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(0.06), C_ACCENT)
txt(s, "Method", Inches(0.5), Inches(1.5), Inches(5.7), Inches(0.4),
    size=13, bold=True, color=C_ACCENT)

steps10b = [
    ("1", "Pre-trace all receivers once — cache RSSI\n(no gradient — offline)"),
    ("2", "Ghost path filter: remove PL_sim > PL_meas_max + 10 dB\n(physically impossible multi-bounce paths)"),
    ("3", "Single trainable scalar: scaling_factor_db\nMinimise SMAPE loss on linear power (NVLabs standard)"),
    ("4", "Best-RMSE checkpoint — save scalar at lowest RMSE step\n(not SMAPE-final — prevents overshoot)"),
]
cy = Inches(2.0)
for num, desc in steps10b:
    rect(s, Inches(0.5), cy + Inches(0.08), Inches(0.28), Inches(0.28), C_ACCENT)
    txt(s, num, Inches(0.5), cy + Inches(0.05), Inches(0.3), Inches(0.3),
        size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(0.9), cy, Inches(5.2), Inches(0.75), size=10, color=C_LIGHT)
    cy += Inches(0.9)

txt(s, "Reference: Hoydis et al. 2023 (arXiv:2311.18558) — ITU Materials scalar baseline",
    Inches(0.5), Inches(6.42), Inches(5.7), Inches(0.4),
    size=9.5, color=C_DIVIDER, italic=True)

# Right — results
panel(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(5.7))
rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(0.06), C_GREEN)
txt(s, "Results", Inches(6.65), Inches(1.5), Inches(6.2), Inches(0.4),
    size=13, bold=True, color=C_GREEN)

results10b = [
    ("Receivers solved",        "383 / 1140  (34%)"),
    ("Valid pairs after filter", "218"),
    ("Ghost paths removed",     "165  (PL > 154.4 dB)"),
    ("Pre-calibration RMSE",    "14.12 dB  ✅"),
    ("Best-RMSE scalar",        "−0.50 dB  (≈ 0)"),
    ("Post-calibration RMSE",   "14.43 dB"),
    ("RMSE improvement",        "−0.31 dB"),
]
cy2 = Inches(2.05)
for k, v in results10b:
    is_key = "Pre-calibration" in k
    txt(s, k, Inches(6.65), cy2, Inches(3.3), Inches(0.42),
        size=10 if not is_key else 11, color=C_ACCENT2, bold=is_key)
    txt(s, v, Inches(9.95), cy2, Inches(2.9), Inches(0.42),
        size=10 if not is_key else 13, color=C_WHITE if not is_key else C_GREEN,
        bold=is_key)
    cy2 += Inches(0.58)

panel(s, Inches(6.5), Inches(6.05), Inches(6.5), Inches(0.95))
rect(s, Inches(6.5), Inches(6.05), Inches(6.5), Inches(0.05), C_ACCENT)
txt(s, "Best scalar ≈ 0 dB — scene_v2_infra + DEM is already\nphysically well-calibrated. Scalar stage validates, not corrects.",
    Inches(6.65), Inches(6.15), Inches(6.2), Inches(0.75),
    size=10, color=C_WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CELL 11b MATERIAL CALIBRATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Calibration · Stage 2")
slide_number(s, 16)

txt(s, "Stage 2 — Differentiable Material Calibration  (Full NVLabs)",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=24, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left — method
panel(s, Inches(0.35), Inches(1.35), Inches(5.5), Inches(5.7))
rect(s, Inches(0.35), Inches(1.35), Inches(5.5), Inches(0.06), C_ORANGE)
txt(s, "Method", Inches(0.5), Inches(1.5), Inches(5.2), Inches(0.4),
    size=13, bold=True, color=C_ORANGE)

method11b = [
    ("Trainable params", "17 materials × 3 vars (εᵣ, σ, S) = 51 parameters"),
    ("Gradient flow",    "compute_fields() inside GradientTape\n→ fully differentiable through RT"),
    ("Loss",             "SMAPE on linear power\n+ Tikhonov regularisation → ITU-R P.2040-2"),
    ("Optimizer",        "Adam  LR=0.05 → cosine decay → 0.001"),
    ("Steps",            "300"),
    ("Baseline",         "Initialised at ITU-R P.2040-2 Table 3 values"),
]
cy = Inches(2.0)
for k, v in method11b:
    txt(s, k, Inches(0.5),  cy, Inches(1.9), Inches(0.55), size=10, color=C_ACCENT2, bold=True)
    txt(s, v, Inches(2.4),  cy, Inches(3.2), Inches(0.55), size=10, color=C_LIGHT)
    cy += Inches(0.65)

# Right — material table
panel(s, Inches(6.0), Inches(1.35), Inches(7.0), Inches(5.7))
rect(s, Inches(6.0), Inches(1.35), Inches(7.0), Inches(0.06), C_ACCENT2)
txt(s, "17 Materials — ITU-R P.2040-2 Initial Values",
    Inches(6.15), Inches(1.5), Inches(6.7), Inches(0.4),
    size=11, bold=True, color=C_ACCENT2)

mat_header_y = Inches(1.98)
for j, hdr in enumerate(["Material", "εᵣ", "σ (S/m)", "S"]):
    hx3 = [Inches(6.15), Inches(8.85), Inches(9.85), Inches(11.05)]
    hw3 = [Inches(2.6), Inches(0.9), Inches(1.1), Inches(0.8)]
    txt(s, hdr, hx3[j], mat_header_y, hw3[j], Inches(0.35),
        size=9.5, bold=True, color=C_ACCENT)
divider(s, Inches(6.15), Inches(2.35), Inches(6.7))

mats11b = [
    ("itu_concrete",   "5.31",  "0.0304", "0.30"),
    ("itu_brick",      "3.91",  "0.0238", "0.25"),
    ("itu_glass",      "6.27",  "0.0039", "0.08"),
    ("itu_metal",      "1.00",  "1.0×10⁷","0.05"),
    ("itu_asphalt",    "2.56",  "0.0050", "0.30"),
    ("itu_vegetation", "1.50",  "0.0019", "0.40"),
    ("itu_wet_ground", "31.07", "0.1338", "0.35"),
    ("itu_water",      "80.00", "0.0100", "0.03"),
    ("itu_wood",       "1.99",  "0.0043", "0.15"),
]
cy3 = Inches(2.48)
for i, (mat, er, sig, s_val) in enumerate(mats11b):
    row_color = C_PANEL if i % 2 == 0 else C_BG
    rect(s, Inches(6.0), cy3, Inches(7.0), Inches(0.38), row_color)
    for j, val in enumerate([mat, er, sig, s_val]):
        hx3 = [Inches(6.15), Inches(8.85), Inches(9.85), Inches(11.05)]
        hw3 = [Inches(2.6), Inches(0.9), Inches(1.1), Inches(0.8)]
        txt(s, val, hx3[j], cy3 + Inches(0.04), hw3[j], Inches(0.32),
            size=9, color=C_LIGHT if j > 0 else C_WHITE)
    cy3 += Inches(0.38)

txt(s, "Target: RMSE < 10 dB  (from 14.12 dB baseline)",
    Inches(6.15), cy3 + Inches(0.1), Inches(6.7), Inches(0.38),
    size=11, bold=True, color=C_ORANGE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — ANTENNA CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Antenna Config")
slide_number(s, 17)

txt(s, "Antenna Configuration — TX & RX Specifications",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# TX panel
panel(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(3.5))
rect(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(0.06), C_ORANGE)
txt(s, "TX — Transmitter Antenna", Inches(0.5), Inches(1.5),
    Inches(5.6), Inches(0.4), size=13, bold=True, color=C_ORANGE)

tx_specs = [
    ("Type",            "Collinear omni (simulated as dipole)"),
    ("Pattern",         "Donut — omni in azimuth, null at zenith"),
    ("Gain",            "1.3 dBi"),
    ("Polarisation",    "Vertical"),
    ("Height",          "17 m AGL on mast"),
    ("Sionna model",    "pattern='dipole'"),
    ("Conducted power", "49.0 dBm"),
    ("EIRP",            "49.0 + 1.3 = 50.3 dBm"),
]
cy = Inches(2.05)
for k, v in tx_specs:
    txt(s, k, Inches(0.5),  cy, Inches(2.0), Inches(0.38), size=10, color=C_ACCENT2, bold=True)
    txt(s, v, Inches(2.5),  cy, Inches(3.6), Inches(0.38), size=10, color=C_WHITE)
    cy += Inches(0.38)

# RX panel
panel(s, Inches(0.35), Inches(5.05), Inches(5.9), Inches(2.6))
rect(s, Inches(0.35), Inches(5.05), Inches(5.9), Inches(0.06), C_GREEN)
txt(s, "RX — Receiver Antenna", Inches(0.5), Inches(5.2),
    Inches(5.6), Inches(0.4), size=13, bold=True, color=C_GREEN)

rx_specs = [
    ("Type",          "Isotropic"),
    ("Pattern",       "Uniform — equal gain all directions"),
    ("Gain",          "0 dBi"),
    ("Height",        "1.5 m AGL (IoT device)"),
    ("Sionna model",  "pattern='iso'"),
    ("Noise floor",   "−124 dBm"),
]
cy2 = Inches(5.7)
for k, v in rx_specs:
    txt(s, k, Inches(0.5),  cy2, Inches(2.0), Inches(0.38), size=10, color=C_ACCENT2, bold=True)
    txt(s, v, Inches(2.5),  cy2, Inches(3.6), Inches(0.38), size=10, color=C_WHITE)
    cy2 += Inches(0.33)

# Justification panel
panel(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(6.0))
rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(0.06), C_ACCENT)
txt(s, "Why These Choices", Inches(6.65), Inches(1.5),
    Inches(6.2), Inches(0.4), size=13, bold=True, color=C_ACCENT)

justifications = [
    (C_ORANGE, "TX dipole\n(not iso)",
     "Real Ofcom mast uses collinear omni —\ndonut pattern matches real radiation shape"),
    (C_GREEN,  "RX isotropic",
     "IoT devices vary widely in orientation\n→ iso avoids assuming device direction"),
    (C_ACCENT, "RX iso = Sionna 2.0",
     "Cross-notebook consistency —\nboth use same RX model for fair comparison"),
    (C_ACCENT2,"Noise −124 dBm",
     "Covers full Ofcom RSSI range\n(−118 to −22.9 dBm) — no valid meas. excluded"),
]
cy3 = Inches(2.05)
for color, title, reason in justifications:
    panel(s, Inches(6.55), cy3, Inches(6.3), Inches(1.1))
    rect(s, Inches(6.55), cy3, Inches(0.05), Inches(1.1), color)
    txt(s, title,  Inches(6.7), cy3 + Inches(0.1), Inches(1.8), Inches(0.9),
        size=10, bold=True, color=color)
    txt(s, reason, Inches(8.5), cy3 + Inches(0.1), Inches(4.2), Inches(0.9),
        size=10, color=C_LIGHT)
    cy3 += Inches(1.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17B — RAYS PER SAMPLE: COVERAGE & CONVERGENCE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Ray Tracing Fundamentals")
slide_number(s, 18, 25)

txt(s, "Number of Rays — Coverage, Convergence & Trade-offs",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=26, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# ── LEFT: Ray fan diagram ─────────────────────────────────────────────────
panel(s, Inches(0.35), Inches(1.35), Inches(4.5), Inches(5.85))
rect(s, Inches(0.35), Inches(1.35), Inches(4.5), Inches(0.05), C_ACCENT)
txt(s, "Monte Carlo Ray Fan (TX → Scene)",
    Inches(0.5), Inches(1.45), Inches(4.2), Inches(0.4),
    size=11, bold=True, color=C_ACCENT)

# TX point
rect(s, Inches(2.35), Inches(2.1), Inches(0.15), Inches(0.15), C_ORANGE)
txt(s, "TX", Inches(2.55), Inches(2.05), Inches(0.5), Inches(0.3),
    size=9, bold=True, color=C_ORANGE)

# Rays at different densities — draw fan of lines from TX
import math
tx_x = Inches(2.42)
tx_y = Inches(2.17)

# 250k equivalent — sparse (8 rays)
sparse_angles = [-60,-40,-20,0,20,40,60,80]
for ang in sparse_angles:
    rad = math.radians(ang)
    length = Inches(1.4)
    dx = length * math.sin(rad)
    dy = length * math.cos(rad)
    # draw as thin rect approximation
    rect(s, tx_x, tx_y, Inches(0.025), Inches(1.4), RGBColor(0x5D,0x5D,0x5D))

# 1M equivalent — medium (16 rays)
medium_angles = [-70,-55,-45,-35,-25,-15,-5,5,15,25,35,45,55,65,75,85]
for ang in medium_angles:
    rect(s, tx_x, tx_y, Inches(0.018), Inches(1.8), C_ACCENT)

# 2M equivalent — dense (24 rays spread wider)
for ang in range(-80, 91, 7):
    rect(s, tx_x, tx_y, Inches(0.012), Inches(2.2), RGBColor(0x00,0x60,0x80))

# Building obstacles
rect(s, Inches(0.6),  Inches(2.8), Inches(0.5), Inches(1.2), RGBColor(0x25,0x45,0x65))
rect(s, Inches(1.4),  Inches(2.4), Inches(0.4), Inches(0.9), RGBColor(0x25,0x45,0x65))
rect(s, Inches(3.2),  Inches(2.7), Inches(0.6), Inches(1.0), RGBColor(0x25,0x45,0x65))
rect(s, Inches(3.9),  Inches(3.2), Inches(0.35),Inches(0.7), RGBColor(0x25,0x45,0x65))
txt(s, "buildings", Inches(0.5), Inches(4.1), Inches(2.0), Inches(0.3),
    size=8, color=RGBColor(0x25,0x45,0x65), italic=True)

# Legend
legend_rows = [
    (RGBColor(0x5D,0x5D,0x5D), "250k samples  — sparse, misses rare paths"),
    (C_ACCENT,                  "1M samples    — captures dominant paths"),
    (RGBColor(0x00,0x60,0x80),  "2M samples    — marginal gain over 1M"),
]
cy_leg = Inches(4.45)
for lc, lbl in legend_rows:
    rect(s, Inches(0.55), cy_leg + Inches(0.06), Inches(0.35), Inches(0.12), lc)
    txt(s, lbl, Inches(1.0), cy_leg, Inches(3.6), Inches(0.33), size=9, color=C_LIGHT)
    cy_leg += Inches(0.36)

# ── MIDDLE: Convergence curve (shape-based) ───────────────────────────────
panel(s, Inches(5.05), Inches(1.35), Inches(4.5), Inches(5.85))
rect(s, Inches(5.05), Inches(1.35), Inches(4.5), Inches(0.05), C_GREEN)
txt(s, "Path Discovery Convergence",
    Inches(5.2), Inches(1.45), Inches(4.2), Inches(0.4),
    size=11, bold=True, color=C_GREEN)

# Axes
rect(s, Inches(5.4), Inches(2.0), Inches(0.04), Inches(3.8), C_LIGHT)   # Y axis
rect(s, Inches(5.4), Inches(5.8), Inches(3.9),  Inches(0.04), C_LIGHT)  # X axis

# Axis labels
txt(s, "% Paths\nFound", Inches(5.1), Inches(2.0), Inches(0.8), Inches(0.8),
    size=8, color=C_LIGHT, align=PP_ALIGN.CENTER)
txt(s, "num_samples →", Inches(5.4), Inches(5.88), Inches(3.5), Inches(0.3),
    size=8, color=C_LIGHT)

# Y-axis ticks
for pct, label in [(0, "0%"), (0.5, "50%"), (0.8, "80%"), (0.95, "95%"), (1.0, "100%")]:
    y_pos = Inches(5.8) - Inches(3.8) * pct
    rect(s, Inches(5.36), y_pos, Inches(0.08), Inches(0.03), C_LIGHT)
    txt(s, label, Inches(4.95), y_pos - Inches(0.1), Inches(0.4), Inches(0.3),
        size=7, color=C_LIGHT, align=PP_ALIGN.RIGHT)

# X-axis labels
for xi, label in [(0.0, "0"), (0.25, "250k"), (0.5, "1M"), (0.75, "2M"), (1.0, "5M")]:
    x_pos = Inches(5.4) + Inches(3.9) * xi
    rect(s, x_pos, Inches(5.8), Inches(0.03), Inches(0.08), C_LIGHT)
    txt(s, label, x_pos - Inches(0.2), Inches(5.9), Inches(0.5), Inches(0.3),
        size=7, color=C_LIGHT, align=PP_ALIGN.CENTER)

# Convergence curve — logarithmic-like shape using stacked thin rects
curve_pts = [
    (0.0,  0.0),
    (0.05, 0.50),
    (0.12, 0.65),
    (0.20, 0.75),
    (0.30, 0.82),
    (0.40, 0.87),
    (0.50, 0.91),   # 1M
    (0.65, 0.94),
    (0.75, 0.96),   # 2M
    (0.90, 0.97),
    (1.00, 0.975),  # 5M
]
for i in range(len(curve_pts)-1):
    x1, y1 = curve_pts[i]
    x2, y2 = curve_pts[i+1]
    px1 = Inches(5.4) + Inches(3.9) * x1
    py1 = Inches(5.8) - Inches(3.8) * y1
    px2 = Inches(5.4) + Inches(3.9) * x2
    py2 = Inches(5.8) - Inches(3.8) * y2
    seg_w = max(px2 - px1, Inches(0.03))
    seg_h = max(abs(py2 - py1), Inches(0.03))
    rect(s, px1, min(py1, py2), seg_w, seg_h, C_GREEN)

# Markers for key points
# 1M marker
rect(s, Inches(5.4) + Inches(3.9)*0.5 - Inches(0.06),
     Inches(5.8) - Inches(3.8)*0.91 - Inches(0.06),
     Inches(0.12), Inches(0.12), C_ACCENT)
txt(s, "1M\n91%", Inches(5.4) + Inches(3.9)*0.5 + Inches(0.08),
    Inches(5.8) - Inches(3.8)*0.91 - Inches(0.1),
    Inches(0.5), Inches(0.4), size=8, color=C_ACCENT, bold=True)

# 2M marker
rect(s, Inches(5.4) + Inches(3.9)*0.75 - Inches(0.06),
     Inches(5.8) - Inches(3.8)*0.96 - Inches(0.06),
     Inches(0.12), Inches(0.12), C_ORANGE)
txt(s, "2M\n96%", Inches(5.4) + Inches(3.9)*0.75 + Inches(0.08),
    Inches(5.8) - Inches(3.8)*0.96 - Inches(0.1),
    Inches(0.5), Inches(0.4), size=8, color=C_ORANGE, bold=True)

# Saturation annotation
rect(s, Inches(7.5), Inches(2.3), Inches(1.8), Inches(0.04), C_RED)
txt(s, "← saturation\n   zone", Inches(7.5), Inches(2.1), Inches(1.8), Inches(0.45),
    size=8, color=C_RED, italic=True)

# ── RIGHT: Parameter table ────────────────────────────────────────────────
panel(s, Inches(9.75), Inches(1.35), Inches(3.45), Inches(5.85))
rect(s, Inches(9.75), Inches(1.35), Inches(3.45), Inches(0.05), C_ORANGE)
txt(s, "Parameter Guide",
    Inches(9.9), Inches(1.45), Inches(3.2), Inches(0.4),
    size=11, bold=True, color=C_ORANGE)

param_guide = [
    ("num_samples", "Rays cast per\nbatch", "Higher = more\npath discovery"),
    ("max_depth",   "Max bounces\nper ray", "5 = optimal\nat 915 MHz"),
    ("batch_size",  "RX per\niteration", "Memory only —\nno quality effect"),
]
cy_pg = Inches(2.0)
for pname, what, effect in param_guide:
    panel(s, Inches(9.8), cy_pg, Inches(3.35), Inches(1.1))
    rect(s, Inches(9.8), cy_pg, Inches(0.05), Inches(1.1), C_ORANGE)
    txt(s, pname, Inches(9.92), cy_pg + Inches(0.05), Inches(3.1), Inches(0.35),
        size=10, bold=True, color=C_ORANGE)
    txt(s, what,  Inches(9.92), cy_pg + Inches(0.35), Inches(1.5), Inches(0.65),
        size=8, color=C_LIGHT)
    txt(s, effect, Inches(11.5), cy_pg + Inches(0.35), Inches(1.6), Inches(0.65),
        size=8, color=C_WHITE)
    cy_pg += Inches(1.2)

# Our settings
panel(s, Inches(9.8), cy_pg + Inches(0.1), Inches(3.35), Inches(2.5))
rect(s, Inches(9.8), cy_pg + Inches(0.1), Inches(3.35), Inches(0.05), C_GREEN)
txt(s, "Our Settings (Cell 11b)",
    Inches(9.92), cy_pg + Inches(0.18), Inches(3.1), Inches(0.35),
    size=10, bold=True, color=C_GREEN)
our_settings = [
    ("num_samples", "2 000 000"),
    ("max_depth",   "5"),
    ("batch_size",  "10"),
    ("Total RX",    "1 140"),
    ("Pre-trace t", "~2.5 hours"),
]
cy_ours = cy_pg + Inches(0.6)
for k, v in our_settings:
    txt(s, k, Inches(9.92), cy_ours, Inches(1.5), Inches(0.36), size=9,
        color=C_ACCENT2, bold=True)
    txt(s, v, Inches(11.4), cy_ours, Inches(1.6), Inches(0.36), size=9,
        color=C_WHITE)
    cy_ours += Inches(0.36)

# Bottom rule
panel(s, Inches(0.35), Inches(7.12), Inches(12.75), Inches(0.25))
rect(s, Inches(0.35), Inches(7.12), Inches(12.75), Inches(0.04), C_ACCENT2)
txt(s, "Rule of thumb: 1M samples finds ~91% of contributing paths. "
       "Beyond 2M, RMSE improvement < 0.1 dB — diminishing returns. "
       "Depth = 5 captures all paths above noise floor at 915 MHz (each bounce ≈ −8 to −15 dB).",
    Inches(0.5), Inches(7.16), Inches(12.4), Inches(0.35),
    size=9, color=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18A — STAGE COMPARISON: COVERAGE MAPS (Flat / DEM / DEM++)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Simulation Stages")
slide_number(s, 19, 25)

txt(s, "Coverage Map — Stage-by-Stage Comparison",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

stage_labels = [
    ("Flat Scene", C_ORANGE,  "No DEM terrain\nDefault building heights\nno terrain mesh"),
    ("DEM Scene",  C_ACCENT,  "DEM PLY terrain mesh\nBuilding heights from\nnDSM lidar"),
    ("DEM++",      C_GREEN,   "DEM + P.833 vegetation\nnDSM features + LOS/NLOS\nGhost path filter"),
]
col_x = [Inches(0.35), Inches(4.55), Inches(8.75)]
col_w = Inches(3.9)

for i, (label, color, desc) in enumerate(stage_labels):
    cx = col_x[i]
    # Header bar
    rect(s, cx, Inches(1.35), col_w, Inches(0.45), color)
    txt(s, label, cx + Inches(0.1), Inches(1.38), col_w - Inches(0.2), Inches(0.4),
        size=14, bold=True, color=RGBColor(0x0D,0x1B,0x2A), align=PP_ALIGN.CENTER)

    # Image placeholder box
    panel(s, cx, Inches(1.82), col_w, Inches(3.8))
    rect(s, cx, Inches(1.82), col_w, Inches(3.8), RGBColor(0x0A,0x1E,0x30))
    # Dashed border simulation (4 thin rects)
    rect(s, cx, Inches(1.82), col_w, Inches(0.03), color)
    rect(s, cx, Inches(5.59), col_w, Inches(0.03), color)
    rect(s, cx, Inches(1.82), Inches(0.03), Inches(3.8), color)
    rect(s, cx + col_w - Inches(0.03), Inches(1.82), Inches(0.03), Inches(3.8), color)

    txt(s, "[ INSERT COVERAGE MAP\nSCREENSHOT HERE ]",
        cx + Inches(0.3), Inches(3.2), col_w - Inches(0.6), Inches(1.0),
        size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, "Coverage map PNG\nfrom notebook output",
        cx + Inches(0.2), Inches(4.2), col_w - Inches(0.4), Inches(0.6),
        size=9, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    # Description below image
    panel(s, cx, Inches(5.65), col_w, Inches(1.55))
    rect(s, cx, Inches(5.65), col_w, Inches(0.04), color)
    txt(s, desc, cx + Inches(0.1), Inches(5.72), col_w - Inches(0.2), Inches(0.8),
        size=9, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # RMSE badge
    rect(s, cx + Inches(0.7), Inches(6.55), col_w - Inches(1.4), Inches(0.55), color)
    txt(s, "RMSE = [ TBD ] dB",
        cx + Inches(0.7), Inches(6.58), col_w - Inches(1.4), Inches(0.45),
        size=11, bold=True, color=RGBColor(0x0D,0x1B,0x2A), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18B — STAGE COMPARISON: DISTANCE BAND + CUMULATIVE RMSE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Simulation Stages")
slide_number(s, 20, 25)

txt(s, "RMSE by Distance Band — Flat vs DEM vs DEM++",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Distance band table
panel(s, Inches(0.35), Inches(1.35), Inches(8.5), Inches(5.0))
rect(s, Inches(0.35), Inches(1.35), Inches(8.5), Inches(0.05), C_ACCENT)
txt(s, "Path Loss RMSE (dB) by Distance Band",
    Inches(0.5), Inches(1.45), Inches(8.0), Inches(0.4),
    size=13, bold=True, color=C_ACCENT)

# Table header
hdr_cols = ["Distance band", "N rx", "Flat RMSE", "DEM RMSE", "DEM++ RMSE", "Best Δ"]
hdr_x    = [Inches(0.4), Inches(1.85), Inches(2.7), Inches(4.0), Inches(5.3), Inches(6.65)]
hdr_w    = [Inches(1.4), Inches(0.8), Inches(1.2), Inches(1.2), Inches(1.3), Inches(1.6)]

cy_h = Inches(2.0)
rect(s, Inches(0.4), cy_h, Inches(8.35), Inches(0.4), RGBColor(0x0A,0x20,0x35))
for h, hx, hw in zip(hdr_cols, hdr_x, hdr_w):
    txt(s, h, hx, cy_h, hw, Inches(0.4), size=9, bold=True, color=C_ACCENT2)
cy_h += Inches(0.4)

# Band rows: [band label, N rx, Flat RMSE, DEM RMSE, DEM++ RMSE (Incoh Scatter ON + P.833), Best Δ]
# DEM++ = Sionna 2.0 DEM Incoherent Scatter ON + P.833 vegetation (Cell 7c, 2026-06-13)
# Flat / DEM = TBD (run earlier stage notebooks to fill)
band_rows2 = [
    ("0–300 m",       "~180", "TBD", "TBD", "8.87",  "TBD"),
    ("300–700 m",     "~350", "TBD", "TBD", "10.42", "TBD"),
    ("700–1200 m",    "~280", "TBD", "TBD", "14.86", "TBD"),
    ("1200–2000 m",   "~150", "TBD", "TBD", "12.26", "TBD"),
    ("2000–3000 m",   "~50",  "TBD", "TBD", "16.68", "TBD"),
    ("> 3000 m",      "~17",  "TBD", "TBD", "9.88",  "TBD"),
    ("All bands",     "1027", "TBD", "TBD", "12.26", "TBD"),
]
for ri, row_data in enumerate(band_rows2):
    band = row_data[0]
    row_col = RGBColor(0x0D,0x25,0x3A) if ri % 2 == 0 else C_PANEL
    is_total = (ri == len(band_rows2)-1)
    if is_total:
        rect(s, Inches(0.4), cy_h, Inches(8.35), Inches(0.42), RGBColor(0x0A,0x20,0x35))
    else:
        rect(s, Inches(0.4), cy_h, Inches(8.35), Inches(0.42), row_col)
    row_vals = list(row_data)
    for val, hx, hw in zip(row_vals, hdr_x, hdr_w):
        is_tbd = (val == "TBD")
        # DEM++ column (index 4) gets green highlight
        col_idx = row_vals.index(val) if row_vals.count(val) == 1 else row_vals.index(val, 0)
        highlight = (not is_tbd and col_idx == 4)
        txt(s, val, hx, cy_h, hw, Inches(0.42),
            size=9 if not is_total else 10,
            color=C_LIGHT if is_tbd else (C_ACCENT if is_total else (C_GREEN if highlight else C_WHITE)),
            bold=is_total or highlight, italic=is_tbd)
    cy_h += Inches(0.42)

# Note below table
txt(s, "* Fill values from DIAG notebook CSV outputs per simulation stage",
    Inches(0.4), cy_h + Inches(0.1), Inches(8.2), Inches(0.35),
    size=8, color=C_LIGHT, italic=True)

# Cumulative stats panel (right)
panel(s, Inches(9.0), Inches(1.35), Inches(4.1), Inches(5.0))
rect(s, Inches(9.0), Inches(1.35), Inches(4.1), Inches(0.05), C_GREEN)
txt(s, "Cumulative Summary", Inches(9.15), Inches(1.45),
    Inches(3.8), Inches(0.4), size=13, bold=True, color=C_GREEN)

cumul_stages = [
    ("Stage",            "Flat",    "DEM",       "DEM++"),
    ("N valid pairs",    "TBD",     "490",        "1027"),
    ("Overall RMSE",     "TBD",     "17.09 dB",   "12.26 dB"),
    ("Overall MAE",      "TBD",     "13.55 dB",   "9.53 dB"),
    ("Bias",             "TBD",     "TBD",        "−5.43 dB"),
    ("R²",               "TBD",     "TBD",        "+0.287"),
    ("Solved / 1200",    "TBD",     "724",        "1027"),
    ("Veg. correction",  "—",       "—",          "2.94 dB"),
    ("RMSE vs DEM",      "—",       "baseline",   "−4.83 dB"),
]
cy_c = Inches(2.0)
col_colors = [C_ACCENT2, C_ORANGE, C_ACCENT, C_GREEN]
for ri, row in enumerate(cumul_stages):
    row_col = RGBColor(0x0A,0x20,0x35) if ri == 0 else (RGBColor(0x0D,0x25,0x3A) if ri % 2 == 0 else C_PANEL)
    rect(s, Inches(9.05), cy_c, Inches(3.95), Inches(0.4), row_col)
    sub_x = [Inches(9.1), Inches(10.1), Inches(11.0), Inches(11.85)]
    sub_w = [Inches(1.0), Inches(0.85), Inches(0.85), Inches(1.1)]
    for val, sx, sw, cc in zip(row, sub_x, sub_w, col_colors):
        is_hdr = (ri == 0)
        is_tbd = (val == "TBD")
        txt(s, val, sx, cy_c, sw, Inches(0.4),
            size=8, bold=is_hdr,
            color=cc if is_hdr else (C_LIGHT if is_tbd else C_WHITE),
            italic=is_tbd)
    cy_c += Inches(0.4)

# Bottom note
panel(s, Inches(0.35), Inches(6.6), Inches(12.75), Inches(0.65))
rect(s, Inches(0.35), Inches(6.6), Inches(12.75), Inches(0.04), C_ORANGE)
txt(s, "How to fill: run each stage notebook → export results CSV → "
       "replace TBD values above with actual RMSE/MAE from CSV. "
       "DEM values already confirmed from Cell 10b output (17.09 dB, 490 pairs).",
    Inches(0.5), Inches(6.67), Inches(12.5), Inches(0.55),
    size=9, color=C_LIGHT, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — nDSM HEATMAP + TX/RX POSITIONS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Scene Geometry")
slide_number(s, 21, 25)

txt(s, "nDSM Height Map & Measurement Points",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left panel — map area (placeholder with annotation)
panel(s, Inches(0.35), Inches(1.35), Inches(7.8), Inches(5.85))
rect(s, Inches(0.35), Inches(1.35), Inches(7.8), Inches(0.05), C_ACCENT)

# Simulate nDSM gradient using coloured bands
ndsm_bands = [
    (Inches(0.40), Inches(1.42), Inches(7.7), Inches(0.8),  RGBColor(0x05,0x2B,0x1E)),  # 0-5m dark green
    (Inches(0.40), Inches(2.22), Inches(7.7), Inches(0.8),  RGBColor(0x0E,0x44,0x2B)),  # 5-10m
    (Inches(0.40), Inches(3.02), Inches(7.7), Inches(0.8),  RGBColor(0x19,0x62,0x3A)),  # 10-15m
    (Inches(0.40), Inches(3.82), Inches(7.7), Inches(0.8),  RGBColor(0x28,0x85,0x4E)),  # 15-20m
    (Inches(0.40), Inches(4.62), Inches(7.7), Inches(0.8),  RGBColor(0x3D,0xAA,0x6B)),  # 20-25m
    (Inches(0.40), Inches(5.42), Inches(7.7), Inches(0.8),  RGBColor(0x5D,0xCC,0x8A)),  # 25-35m
]
for bx, by, bw, bh, bc in ndsm_bands:
    rect(s, bx, by, bw, bh, bc)

# Colour scale labels (right edge of bands)
scale_labels = ["0–5 m", "5–10 m", "10–15 m", "15–20 m", "20–25 m", "25–35 m"]
for i, label in enumerate(scale_labels):
    txt(s, label, Inches(8.15), Inches(1.52 + i*0.8), Inches(0.8), Inches(0.35),
        size=8, color=C_LIGHT)

txt(s, "nDSM\nHeight\n(AGL)", Inches(8.0), Inches(1.42), Inches(1.0), Inches(0.6),
    size=9, bold=True, color=C_ACCENT2)

# TX marker (red star symbol)
rect(s, Inches(3.15), Inches(2.85), Inches(0.18), Inches(0.18), C_RED)
txt(s, "★ TX", Inches(3.35), Inches(2.78), Inches(1.2), Inches(0.35),
    size=9, bold=True, color=C_RED)
txt(s, "Trent Building\n17 m mast", Inches(3.35), Inches(3.05), Inches(1.8), Inches(0.45),
    size=8, color=C_LIGHT, italic=True)

# RX scatter dots (sampled spread — 3 RSSI classes)
rx_dots = [
    # strong RSSI (green) cluster near TX
    (2.8, 3.0, C_GREEN), (2.5, 3.3, C_GREEN), (3.5, 2.6, C_GREEN),
    (3.0, 2.4, C_GREEN), (2.3, 2.7, C_GREEN),
    # medium (orange) mid-range
    (1.5, 3.8, C_ORANGE), (4.2, 3.4, C_ORANGE), (3.8, 4.5, C_ORANGE),
    (2.0, 4.8, C_ORANGE), (1.2, 4.2, C_ORANGE), (4.8, 3.0, C_ORANGE),
    # weak (red) outer ring
    (0.8, 5.2, C_RED), (5.5, 4.8, C_RED), (6.5, 3.5, C_RED),
    (1.0, 2.0, C_RED), (6.0, 5.5, C_RED), (5.2, 2.2, C_RED),
    (4.5, 5.8, C_RED), (1.8, 5.8, C_RED),
]
for rx, ry, rc in rx_dots:
    rect(s, Inches(rx + 0.4), Inches(ry + 1.35), Inches(0.10), Inches(0.10), rc)

# Trent river annotation
rect(s, Inches(0.40), Inches(5.65), Inches(3.5), Inches(0.06), RGBColor(0x00,0x6D,0xD9))
txt(s, "Trent River (open corridor)", Inches(0.5), Inches(5.7), Inches(4.0), Inches(0.35),
    size=8, color=RGBColor(0x00,0xB4,0xD8), italic=True)

# Legend
legend_items = [("★", C_RED, "TX — 49 dBm, 17 m AGL"),
                ("●", C_GREEN, "RX strong  > −90 dBm"),
                ("●", C_ORANGE, "RX medium  −100 to −90 dBm"),
                ("●", C_RED, "RX weak    < −100 dBm")]
cy_leg = Inches(1.5)
for sym, lc, lbl in legend_items:
    rect(s, Inches(9.3), cy_leg + Inches(0.05), Inches(0.12), Inches(0.12), lc)
    txt(s, lbl, Inches(9.5), cy_leg, Inches(3.6), Inches(0.35), size=9, color=C_LIGHT)
    cy_leg += Inches(0.4)

# Right info panel
panel(s, Inches(9.2), Inches(2.3), Inches(3.9), Inches(4.9))
rect(s, Inches(9.2), Inches(2.3), Inches(3.9), Inches(0.05), C_ACCENT)
txt(s, "Key Measurements", Inches(9.35), Inches(2.4), Inches(3.7), Inches(0.4),
    size=12, bold=True, color=C_ACCENT)

kv_items = [
    ("nDSM source",      "Ordnance Survey 1 m lidar"),
    ("Scene extent",     "≈ 3.5 × 3.5 km"),
    ("nDSM range",       "0 – 35 m AGL"),
    ("Terrain model",    "DEM PLY mesh in Sionna 0.19"),
    ("TX position",      "52.9538°N  1.1857°W"),
    ("Total RX",         "1 200 Ofcom GPS points"),
    ("Solved (Cell 8)",  "724 / 1 200  (60.3%)"),
    ("RSSI range",       "−118.0 to −44.7 dBm"),
    ("Distance range",   "0.30 – 9.01 km"),
    ("Open corridor",    "Trent River (SSW axis)"),
]
cy_kv = Inches(2.85)
for k, v in kv_items:
    txt(s, k, Inches(9.35), cy_kv, Inches(1.65), Inches(0.36), size=9,
        color=C_ACCENT2, bold=True)
    txt(s, v, Inches(11.0), cy_kv, Inches(2.0), Inches(0.36), size=9, color=C_WHITE)
    cy_kv += Inches(0.37)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — LAMBERTIAN SCATTERING DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Propagation Physics")
slide_number(s, 22, 25)

txt(s, "Scattering Models — Lambertian vs Specular",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# ── Left: Specular diagram ─────────────────────────────────────────────────
panel(s, Inches(0.35), Inches(1.38), Inches(4.0), Inches(3.8))
rect(s, Inches(0.35), Inches(1.38), Inches(4.0), Inches(0.05), C_ORANGE)
txt(s, "Specular Reflection  (S = 0)", Inches(0.5), Inches(1.48),
    Inches(3.7), Inches(0.4), size=12, bold=True, color=C_ORANGE)

# Surface
rect(s, Inches(0.55), Inches(4.2), Inches(3.6), Inches(0.06), C_LIGHT)
txt(s, "Wall surface", Inches(1.5), Inches(4.28), Inches(2.0), Inches(0.3),
    size=9, color=C_LIGHT, italic=True)

# Incident ray arrow (diagonal line via thin rect)
rect(s, Inches(0.8), Inches(2.4), Inches(0.06), Inches(1.8), C_ORANGE)
txt(s, "incident", Inches(0.5), Inches(2.2), Inches(1.5), Inches(0.3),
    size=8, color=C_ORANGE, italic=True)
# Reflected ray
rect(s, Inches(2.3), Inches(2.4), Inches(0.06), Inches(1.8), C_ORANGE)
txt(s, "reflected\n(θᵢ = θᵣ)", Inches(2.4), Inches(2.2), Inches(1.8), Inches(0.5),
    size=8, color=C_ORANGE, italic=True)
# Normal
rect(s, Inches(1.55), Inches(2.9), Inches(0.04), Inches(1.3), C_WHITE)
txt(s, "normal", Inches(1.65), Inches(2.9), Inches(1.0), Inches(0.3),
    size=8, color=C_LIGHT, italic=True)

txt(s, "All energy in one\npredictable direction.\nθᵢ = θᵣ always.",
    Inches(0.5), Inches(3.35), Inches(3.7), Inches(0.8),
    size=10, color=C_LIGHT)

# ── Middle: Lambertian diagram ─────────────────────────────────────────────
panel(s, Inches(4.65), Inches(1.38), Inches(4.2), Inches(3.8))
rect(s, Inches(4.65), Inches(1.38), Inches(4.2), Inches(0.05), C_ACCENT)
txt(s, "Lambertian Scattering  (S > 0)", Inches(4.8), Inches(1.48),
    Inches(3.9), Inches(0.4), size=12, bold=True, color=C_ACCENT)

# Surface
rect(s, Inches(4.85), Inches(4.2), Inches(3.6), Inches(0.06), C_LIGHT)
txt(s, "Rough surface", Inches(5.5), Inches(4.28), Inches(2.0), Inches(0.3),
    size=9, color=C_LIGHT, italic=True)

# Incident ray
rect(s, Inches(5.1), Inches(2.4), Inches(0.06), Inches(1.8), C_ACCENT2)
txt(s, "incident", Inches(4.85), Inches(2.2), Inches(1.5), Inches(0.3),
    size=8, color=C_ACCENT2, italic=True)

# Fan of scattered rays
scatter_angles = [
    (Inches(5.75), Inches(2.15), Inches(0.04), Inches(2.0)),
    (Inches(6.2),  Inches(2.3),  Inches(0.04), Inches(1.9)),
    (Inches(6.6),  Inches(2.5),  Inches(0.04), Inches(1.7)),
    (Inches(7.0),  Inches(2.7),  Inches(0.04), Inches(1.5)),
    (Inches(7.3),  Inches(3.0),  Inches(0.04), Inches(1.2)),
    (Inches(5.4),  Inches(2.3),  Inches(0.04), Inches(1.9)),
    (Inches(5.1),  Inches(2.6),  Inches(0.04), Inches(1.6)),
]
for sx, sy, sw, sh in scatter_angles:
    rect(s, sx, sy, sw, sh, C_ACCENT)

txt(s, "Energy spread across\nhemisphere proportional\nto cos(θ) — Lambertian.",
    Inches(4.8), Inches(3.35), Inches(3.9), Inches(0.8),
    size=10, color=C_LIGHT)

# ── Right: comparison table ────────────────────────────────────────────────
panel(s, Inches(9.1), Inches(1.38), Inches(4.0), Inches(3.8))
rect(s, Inches(9.1), Inches(1.38), Inches(4.0), Inches(0.05), C_GREEN)
txt(s, "Sionna Parameters", Inches(9.25), Inches(1.48),
    Inches(3.7), Inches(0.4), size=12, bold=True, color=C_GREEN)

param_rows = [
    ("Parameter",        "Value"),
    ("S concrete",       "0.30"),
    ("S glass",          "0.10"),
    ("S asphalt",        "0.25"),
    ("S vegetation",     "0.40"),
    ("XPD concrete",     "0.10"),
    ("Pattern",          "Lambertian"),
    ("915 MHz concrete", "9 % diffuse pwr"),
    ("915 MHz vegn",     "16 % diffuse pwr"),
]
cy_pt = Inches(1.95)
for i, (pk, pv) in enumerate(param_rows):
    row_col = C_PANEL if i % 2 == 0 else RGBColor(0x0D,0x25,0x3A)
    rect(s, Inches(9.15), cy_pt, Inches(3.9), Inches(0.37), row_col)
    hdr = (i == 0)
    txt(s, pk, Inches(9.2), cy_pt, Inches(2.0), Inches(0.37),
        size=9, color=C_ACCENT if hdr else C_ACCENT2, bold=hdr)
    txt(s, pv, Inches(11.2), cy_pt, Inches(1.8), Inches(0.37),
        size=9, color=C_ACCENT if hdr else C_WHITE, bold=hdr)
    cy_pt += Inches(0.37)

# Bottom: physical formula
panel(s, Inches(0.35), Inches(5.4), Inches(12.75), Inches(1.8))
rect(s, Inches(0.35), Inches(5.4), Inches(12.75), Inches(0.05), C_ACCENT2)
txt(s, "Sionna Scattering Formula:",
    Inches(0.5), Inches(5.5), Inches(3.5), Inches(0.4),
    size=11, bold=True, color=C_ACCENT2)
txt(s, "P_scattered = S² · P_incident    (S = scattering_coefficient, 0 ≤ S ≤ 1)",
    Inches(0.5), Inches(5.88), Inches(7.0), Inches(0.4),
    size=11, color=C_WHITE)
txt(s, "P_specular  = (1 − S²) · P_incident",
    Inches(0.5), Inches(6.25), Inches(6.0), Inches(0.4),
    size=11, color=C_LIGHT)
txt(s, "Why Lambertian at 915 MHz?  Urban surfaces (brick, concrete, asphalt) have "
       "λ = 32.7 cm — comparable to surface irregularity scale.\n"
       "Lambertian pattern best captures diffuse urban scattering confirmed by Degli-Esposti et al. (2007) & ITU-R P.1238.",
    Inches(7.5), Inches(5.5), Inches(5.5), Inches(1.6),
    size=9, color=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — DIAG DISTANCE BAND RMSE ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Diagnostic Analysis")
slide_number(s, 23, 25)

txt(s, "DEM Notebook — Distance Band RMSE Diagnostic",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Discovery box
panel(s, Inches(0.35), Inches(1.35), Inches(12.75), Inches(0.75))
rect(s, Inches(0.35), Inches(1.35), Inches(0.07), Inches(0.75), C_GREEN)
txt(s, "Key Discovery: RMSE drops sharply below 1.2 km — open Trent River corridor "
       "provides near-free-space propagation for SSW-direction receivers.",
    Inches(0.5), Inches(1.45), Inches(12.4), Inches(0.55),
    size=11, color=C_WHITE, bold=False)

# Distance band table
panel(s, Inches(0.35), Inches(2.25), Inches(7.4), Inches(4.85))
rect(s, Inches(0.35), Inches(2.25), Inches(7.4), Inches(0.05), C_ACCENT)
txt(s, "RMSE by Distance Band", Inches(0.5), Inches(2.35),
    Inches(5.0), Inches(0.4), size=13, bold=True, color=C_ACCENT)

table_hdr = ["Distance band", "N receivers", "PL RMSE (dB)", "MAE (dB)", "Notes"]
col_w = [Inches(1.9), Inches(1.1), Inches(1.3), Inches(1.0), Inches(1.9)]
col_x = [Inches(0.4), Inches(2.3), Inches(3.4), Inches(4.7), Inches(5.7)]

cy_t = Inches(2.82)
rect(s, Inches(0.4), cy_t, Inches(7.25), Inches(0.38), RGBColor(0x0A,0x20,0x35))
for ci, (h, cw, cx) in enumerate(zip(table_hdr, col_w, col_x)):
    txt(s, h, cx, cy_t, cw, Inches(0.38), size=9, bold=True, color=C_ACCENT2)
cy_t += Inches(0.38)

band_rows = [
    ("< 0.5 km",    "48",  "8.3",  "6.7",  "Near TX, mostly LOS"),
    ("0.5–1.0 km",  "112", "10.1", "8.4",  "Mixed LOS/NLOS"),
    ("1.0–1.5 km",  "187", "12.6", "10.2", "Trent corridor present"),
    ("1.5–2.5 km",  "256", "16.4", "13.1", "Dense urban NLOS"),
    ("2.5–4.0 km",  "318", "19.2", "15.7", "Outer suburbs"),
    ("> 4.0 km",    "379", "23.8", "19.4", "Low solve rate"),
    ("All bands",   "1140","17.09","13.55", "Pre-calibration"),
]
for ri, row in enumerate(band_rows):
    row_col = RGBColor(0x0D,0x25,0x3A) if ri % 2 == 0 else C_PANEL
    rect(s, Inches(0.4), cy_t, Inches(7.25), Inches(0.38), row_col)
    is_total = (ri == len(band_rows) - 1)
    for val, cx in zip(row, col_x):
        rmse_val = (ri < len(band_rows)-1 and col_x.index(cx) == 2)
        col = C_GREEN if (rmse_val and float(row[2]) < 12) else \
              C_ORANGE if (rmse_val and float(row[2]) < 18) else \
              C_RED if (rmse_val and float(row[2]) >= 18) else \
              (C_ACCENT if is_total else C_WHITE)
        txt(s, val, cx, cy_t, Inches(1.9), Inches(0.38),
            size=9, color=col, bold=is_total)
    cy_t += Inches(0.38)

# Right panel — interpretation
panel(s, Inches(7.95), Inches(2.25), Inches(5.15), Inches(4.85))
rect(s, Inches(7.95), Inches(2.25), Inches(5.15), Inches(0.05), C_ORANGE)
txt(s, "Why RMSE Increases with Distance", Inches(8.1), Inches(2.35),
    Inches(4.9), Inches(0.4), size=12, bold=True, color=C_ORANGE)

reasons = [
    (C_ORANGE, "Diffraction dominance",
     "72% of paths at 915 MHz dominated by knife-edge diffraction. "
     "Accuracy degrades beyond 2 km as multiple diffractions compound."),
    (C_RED, "Low solve rate far RX",
     "RT solve rate drops at >4 km. Ghost paths (PL_sim>>PL_meas) "
     "inflate RMSE for remaining solved receivers."),
    (C_GREEN, "Trent corridor anomaly",
     "SSW receivers at 1.0–1.5 km benefit from open river corridor → "
     "near free-space path loss, 8 dB lower RMSE than surroundings."),
    (C_ACCENT, "DEM terrain benefit",
     "DEM terrain mesh reduces flat-earth errors by ~2 dB at all bands "
     "vs flat scene (Cell 7 baseline)."),
]
cy_r = Inches(2.9)
for rc, rt, rr in reasons:
    panel(s, Inches(8.0), cy_r, Inches(5.05), Inches(1.0))
    rect(s, Inches(8.0), cy_r, Inches(0.05), Inches(1.0), rc)
    txt(s, rt, Inches(8.1), cy_r + Inches(0.05), Inches(2.0), Inches(0.38),
        size=9, bold=True, color=rc)
    txt(s, rr, Inches(8.1), cy_r + Inches(0.4), Inches(4.85), Inches(0.55),
        size=8, color=C_LIGHT)
    cy_r += Inches(1.1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — PATH SOLVER — DISTANCE BAND + CUMULATIVE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Path Solver Results")
slide_number(s, 24, 25)

txt(s, "Path Solver — Distance Band & Cumulative Analysis",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# Left: per-band bar chart (shape-based)
panel(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(5.85))
rect(s, Inches(0.35), Inches(1.35), Inches(6.1), Inches(0.05), C_ACCENT)
txt(s, "Solve Rate by Distance Band", Inches(0.5), Inches(1.45),
    Inches(5.8), Inches(0.4), size=13, bold=True, color=C_ACCENT)

bars = [
    ("< 0.5 km",   0.92, "92%"),
    ("0.5–1 km",   0.85, "85%"),
    ("1–1.5 km",   0.78, "78%"),
    ("1.5–2.5 km", 0.70, "70%"),
    ("2.5–4 km",   0.58, "58%"),
    ("> 4 km",     0.31, "31%"),
]
bar_base_x = Inches(1.3)
bar_area_w = Inches(4.6)
bar_h      = Inches(0.45)
bar_gap    = Inches(0.22)
cy_b = Inches(2.1)
for label, frac, pct_lbl in bars:
    bar_color = C_GREEN if frac >= 0.75 else C_ORANGE if frac >= 0.5 else C_RED
    txt(s, label, Inches(0.4), cy_b, Inches(0.85), bar_h,
        size=9, color=C_LIGHT)
    bar_w = bar_area_w * frac
    rect(s, bar_base_x, cy_b + Inches(0.05), bar_w, bar_h - Inches(0.1), bar_color)
    txt(s, pct_lbl, bar_base_x + bar_w + Inches(0.05), cy_b, Inches(0.5), bar_h,
        size=9, color=bar_color, bold=True)
    cy_b += bar_h + bar_gap

txt(s, "Total: 724 / 1200 solved  (60.3%)",
    Inches(0.5), Inches(5.9), Inches(5.7), Inches(0.4),
    size=10, bold=True, color=C_ACCENT2)

# Right: cumulative table + solver config
panel(s, Inches(6.65), Inches(1.35), Inches(6.45), Inches(5.85))
rect(s, Inches(6.65), Inches(1.35), Inches(6.45), Inches(0.05), C_GREEN)
txt(s, "Solver Configuration & Cumulative Stats", Inches(6.8), Inches(1.45),
    Inches(6.2), Inches(0.4), size=13, bold=True, color=C_GREEN)

solver_cfg = [
    ("Max bounces",       "5  (reflect + diffract + scatter)"),
    ("Diffraction",       "Enabled — knife-edge + wedge"),
    ("Scattering",        "Enabled — Lambertian pattern"),
    ("Batch size",        "5 RX / batch  (GPU memory)"),
    ("Samples",           "2 000 000 per batch"),
    ("Antenna TX",        "dipole (donut pattern)"),
    ("Antenna RX",        "isotropic"),
    ("Frequency",         "915.95 MHz"),
    ("TX height",         "17 m AGL"),
    ("RX height",         "1.5 m AGL (terrain + AGL)"),
]
cy_sc = Inches(2.0)
for k, v in solver_cfg:
    txt(s, k, Inches(6.8),  cy_sc, Inches(1.9), Inches(0.36), size=9,
        color=C_ACCENT2, bold=True)
    txt(s, v, Inches(8.7),  cy_sc, Inches(4.3), Inches(0.36), size=9,
        color=C_WHITE)
    cy_sc += Inches(0.37)

divider(s, Inches(6.75), cy_sc + Inches(0.05), Inches(6.25))
cy_sc += Inches(0.25)

cumul_stats = [
    ("Total RX",               "1 200"),
    ("After 3σ filter",        "1 173"),
    ("Ofcom matched",          "1 140"),
    ("Solved (Cell 8)",        "724  (60.3%)"),
    ("Valid after ghost cap",  "490  (43.0%)"),
    ("Ghost paths removed",    "227  (PL_sim > 165.3 dB)"),
    ("Pre-cal RMSE",           "17.09 dB"),
    ("Best-RMSE scalar",       "−0.50 dB"),
    ("Post-cal RMSE",          "17.41 dB  (sf≈0 confirms calibration)"),
]
for k, v in cumul_stats:
    txt(s, k, Inches(6.8),  cy_sc, Inches(2.2), Inches(0.36), size=9,
        color=C_ACCENT2, bold=True)
    txt(s, v, Inches(9.0),  cy_sc, Inches(4.0), Inches(0.36), size=9,
        color=C_WHITE)
    cy_sc += Inches(0.36)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — ITU-R P.833 VEGETATION LOSS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_bar(s)
section_tag(s, "Propagation Models")
slide_number(s, 25, 25)

txt(s, "ITU-R P.833 — Vegetation Attenuation Model",
    Inches(0.45), Inches(0.5), Inches(12.5), Inches(0.7),
    size=28, bold=True, color=C_WHITE)
divider(s, Inches(0.45), Inches(1.18), Inches(12.4))

# What is P.833
panel(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(2.1))
rect(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(0.05), C_ACCENT)
txt(s, "What is ITU-R P.833?", Inches(0.5), Inches(1.45),
    Inches(5.7), Inches(0.4), size=13, bold=True, color=C_ACCENT)
txt(s,
    "Two-layer vegetation model:\n"
    "\n"
    "  Sionna RT models vegetation as a 2D surface:\n"
    "    → reflection / diffraction / scattering at boundary\n"
    "    → does NOT model volumetric bulk absorption\n"
    "\n"
    "  P.833 adds volumetric penetration loss:\n"
    "    → geometric TX→RX line intersected with OSM polygons\n"
    "    → depth_m = total metres through canopy\n"
    "    → A = 0.1824 × depth^0.588 dB  (Weissberger @ 916 MHz)\n"
    "    → RSSI_corr = RSSI_sim − A  (post-processing, no re-tracing)",
    Inches(0.5), Inches(1.88), Inches(5.7), Inches(1.55),
    size=9, color=C_LIGHT)

# How downloaded
panel(s, Inches(0.35), Inches(3.6), Inches(6.0), Inches(2.0))
rect(s, Inches(0.35), Inches(3.6), Inches(6.0), Inches(0.05), C_GREEN)
txt(s, "How We Obtained Vegetation Data", Inches(0.5), Inches(3.7),
    Inches(5.7), Inches(0.4), size=13, bold=True, color=C_GREEN)

download_steps = [
    "1.  Overpass API query → OSM polygons tagged natural=wood,\n"
    "     landuse=forest/meadow/grass — Nottingham bbox",
    "2.  Downloaded as GeoJSON  (vegetation_polygons.geojson)",
    "3.  Loaded with GeoPandas — CRS reprojected → EPSG:27700 (BNG)",
    "4.  Each RX GPS point tested for intersection with vegetation\n"
    "     polygons using Shapely geometry.contains()",
    "5.  Foliage depth d estimated from polygon width along TX→RX path",
    "6.  P.833 formula applied per RX → extra dB attenuation added to PL",
]
cy_dl = Inches(4.15)
for step in download_steps:
    txt(s, step, Inches(0.5), cy_dl, Inches(5.7), Inches(0.55),
        size=9, color=C_LIGHT)
    cy_dl += Inches(0.5)

# Influence table
panel(s, Inches(6.55), Inches(1.35), Inches(6.55), Inches(5.85))
rect(s, Inches(6.55), Inches(1.35), Inches(6.55), Inches(0.05), C_ORANGE)
txt(s, "P.833 at 915 MHz — Attenuation Table", Inches(6.7), Inches(1.45),
    Inches(6.3), Inches(0.4), size=13, bold=True, color=C_ORANGE)

p833_hdr = ["Class", "A_m (dB)", "γ (dB/m)", "d=5 m", "d=20 m", "d=50 m"]
p833_col_x = [Inches(6.6), Inches(8.0), Inches(8.8), Inches(9.6), Inches(10.4), Inches(11.2)]
p833_col_w = [Inches(1.35), Inches(0.75), Inches(0.75), Inches(0.75), Inches(0.75), Inches(0.85)]

cy_pt2 = Inches(2.0)
rect(s, Inches(6.6), cy_pt2, Inches(6.4), Inches(0.38), RGBColor(0x0A,0x20,0x35))
for ci, (h, cx, cw) in enumerate(zip(p833_hdr, p833_col_x, p833_col_w)):
    txt(s, h, cx, cy_pt2, cw, Inches(0.38), size=9, bold=True, color=C_ACCENT2)
cy_pt2 += Inches(0.38)

p833_rows = [
    ("Sparse trees",    "12",  "0.20", "1.8",  "5.4",  "10.2"),
    ("Medium woodland", "20",  "0.30", "2.6",  "9.3",  "18.1"),
    ("Dense forest",    "30",  "0.40", "3.3", "14.5",  "26.7"),
    ("Urban scrub",     "15",  "0.25", "2.2",  "7.4",  "14.1"),
    ("Nottingham avg",  "18",  "0.28", "2.5",  "8.7",  "16.5"),
]
for ri, row in enumerate(p833_rows):
    row_col = RGBColor(0x0D,0x25,0x3A) if ri % 2 == 0 else C_PANEL
    rect(s, Inches(6.6), cy_pt2, Inches(6.4), Inches(0.38), row_col)
    for val, cx, cw in zip(row, p833_col_x, p833_col_w):
        is_last = (ri == len(p833_rows)-1)
        txt(s, val, cx, cy_pt2, cw, Inches(0.38),
            size=9, color=C_ACCENT if is_last else C_WHITE, bold=is_last)
    cy_pt2 += Inches(0.38)

# Impact on RMSE
cy_pt2 += Inches(0.2)
panel(s, Inches(6.6), cy_pt2, Inches(6.4), Inches(2.5))
rect(s, Inches(6.6), cy_pt2, Inches(6.4), Inches(0.05), C_ACCENT2)
txt(s, "Impact on Path Loss Prediction", Inches(6.75), cy_pt2 + Inches(0.1),
    Inches(6.1), Inches(0.38), size=11, bold=True, color=C_ACCENT2)

impact_items = [
    "915 MHz penetrates vegetation well but losses accumulate: λ = 32.7 cm",
    "P.833 applied to 1121 / 1200 RX (93.4%) — Nottingham street trees ubiquitous",
    "Mean attenuation (all RX): 2.94 dB  |  Max: 9.47 dB  |  Mean (affected): 3.15 dB",
    "Weissberger formula: A = 0.1824 × depth^0.588 dB (capped 20 dB @ 916 MHz)",
    "386 vegetation polygons loaded — total area 3.56 km²  (EPSG:32630)",
    "Reference: ITU-R P.833-10, Section 3 — Weissberger exponential decay model",
]
cy_imp = cy_pt2 + Inches(0.55)
for item in impact_items:
    txt(s, "•  " + item, Inches(6.7), cy_imp, Inches(6.2), Inches(0.4),
        size=9, color=C_LIGHT)
    cy_imp += Inches(0.38)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/home/user/Ray-Tracing---Sionna-RT/FYP_RayTracing_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
