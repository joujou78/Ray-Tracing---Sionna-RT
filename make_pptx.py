from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
C_ACCENT   = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_WARN     = RGBColor(0xFF, 0x6B, 0x35)   # orange
C_OK       = RGBColor(0x06, 0xD6, 0xA0)   # green
C_LIGHT    = RGBColor(0xE0, 0xE0, 0xE0)   # light grey text
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD     = RGBColor(0x1A, 0x2E, 0x44)   # slightly lighter navy

W, H = Inches(13.33), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper: add a rectangle ───────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, alpha=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, size=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def heading(slide, text, y=Inches(0.18)):
    txt(slide, text, Inches(0.4), y, Inches(12.5), Inches(0.6),
        size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)

def bullet_box(slide, items, x, y, w, h, title=None, title_color=C_ACCENT,
               item_size=15, title_size=17):
    rect(slide, x, y, w, h, C_CARD)
    cy = y + Inches(0.1)
    if title:
        txt(slide, title, x+Inches(0.15), cy, w-Inches(0.2), Inches(0.35),
            size=title_size, bold=True, color=title_color)
        cy += Inches(0.38)
    for item in items:
        color  = C_WARN  if item.startswith("✗") else \
                 C_OK    if item.startswith("✓") else \
                 C_ACCENT if item.startswith("▶") else C_LIGHT
        txt(slide, item, x+Inches(0.18), cy, w-Inches(0.3), Inches(0.35),
            size=item_size, color=color)
        cy += Inches(0.33)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, Inches(2.8), W, Inches(2.1), C_CARD)

txt(s, "Sionna 0.19.2 Ray-Tracing", Inches(0.6), Inches(1.1),
    Inches(12), Inches(1.0), size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(s, "Nottingham Urban Scene · 3602.5 MHz · 4001 Ofcom Drive-Test Points",
    Inches(0.6), Inches(2.1), Inches(12), Inches(0.6),
    size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Project Milestones · Failures & Fixes · Scene Enhancements",
    Inches(0.6), Inches(3.05), Inches(12), Inches(0.5),
    size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
txt(s, "FYP 2026  ·  Sionna RT + AWS DEM + OSM + Overture Maps + Meta CHM",
    Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
    size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, 0, W, Inches(0.9), C_CARD)
heading(s, "Project Overview")

cols = [
    ("Goal", C_ACCENT, [
        "▶ Simulate urban radio propagation",
        "▶ Nottingham city, 10.7 × 11.2 km",
        "▶ 3602.5 MHz (5G NR n78 band)",
        "▶ Compare vs 4001 Ofcom RSSI points",
        "▶ No bias / correction factors",
    ]),
    ("RF Parameters", C_OK, [
        "TX conducted power: 45.0 dBm",
        "EIRP: 54.0 dBm (9 dBi antenna)",
        "RX extra gain: 18.0 dB (LNA)",
        "Noise floor: −109 dBm (Ofcom)",
        "SITE_CORRECTION_DB = 0.0",
    ]),
    ("Simulation Engine", C_WARN, [
        "Sionna 0.19.2 (TF 2.15 + Mitsuba 3)",
        "SBR path solver (GPU: cuda_ad_rgb)",
        "MAX_DEPTH = 10 bounces",
        "NUM_SAMPLES_CM = 200M rays",
        "NUM_SAMPLES_PS = 10M / batch",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    bullet_box(s, items, Inches(0.3 + i*4.35), Inches(1.0),
               Inches(4.2), Inches(5.8), title=title, title_color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Scene Evolution (checkpoints)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, 0, W, Inches(0.9), C_CARD)
heading(s, "Scene Evolution — Saved Checkpoints")

stages = [
    ("CP 1\nFlat Terrain\n+ OSM Buildings",
     ["95k individual PLY files", "8 merged PLY shapes", "Flat ground z=0", "itu_wet_ground material"],
     C_WARN),
    ("CP 2\nAWS DEM\nTerrain",
     ["72 GeoTIFF tiles (zoom-14)", "500×500 terrain grid", "Z range: −34 → +87 m", "Bilinear interpolation"],
     C_ACCENT),
    ("CP 3\nOverture Maps\nBuilding Heights",
     ["70,765 buildings fetched", "22,558 heights enriched", "Fallback: OSM levels×3.5m", "Default: 8.0 m"],
     C_OK),
    ("CP 4\nMeta CHM\nTree Geometry",
     ["Quadkey zoom-9 tiles", "Tile 031311330 (1–53 m)", "115,497 tree cylinders", "89 MB binary PLY"],
     C_OK),
]

for i, (title, items, color) in enumerate(stages):
    x = Inches(0.25 + i * 3.27)
    rect(s, x, Inches(1.05), Inches(3.1), Inches(5.7), C_CARD)
    # title banner
    rect(s, x, Inches(1.05), Inches(3.1), Inches(0.85), color)
    txt(s, title, x+Inches(0.1), Inches(1.07), Inches(2.9), Inches(0.82),
        size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    cy = Inches(2.05)
    for item in items:
        txt(s, "• " + item, x+Inches(0.15), cy, Inches(2.85), Inches(0.35),
            size=13, color=C_LIGHT)
        cy += Inches(0.36)
    # arrow between stages
    if i < 3:
        txt(s, "→", x+Inches(3.1), Inches(3.55), Inches(0.3), Inches(0.4),
            size=24, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Failures & Fixes
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, 0, W, Inches(0.9), C_CARD)
heading(s, "Key Failures & Fixes")

failures = [
    ("Flat PL-FSPL = +10.3 dB at ALL distances",
     "itu_wet_ground (ε=30) created strong specular bounces dominating all paths",
     "Switched to itu_urban_ground (ε=5, σ=0.5) — absorbs ground reflections"),
    ("Meta CHM HTTP 404 (lat_lon.tif naming)",
     "Tiles use Web Mercator quadkey zoom-9 format, not degree-grid lat/lon",
     "Implemented _latlon_to_quadkey() — Nottingham = tile 031311330"),
    ("OOM: two 65536×65536 CHM tiles = 32 GB RAM",
     "Loading full tiles before mosaic killed the kernel",
     "Windowed reads: clip each tile to scene bbox before loading (~18 MB each)"),
    ("0 trees placed despite valid CHM data",
     "Manual mosaic row-offset bug + 26M-point Python loop (hours to run)",
     "Vectorised numpy meshgrid in Mercator coords, per-strip direct indexing"),
    ("95,185 PLY files crashed scene.xml (12 MB)",
     "Individual bld_XXXXX_wall/roof.ply for each OSM building",
     "Merged by material → 8 PLY files (wall_itu_brick.ply etc.)"),
    ("NaN tag bug: 95% of buildings got wrong material",
     "str(NaN)='nan' is truthy — if off: return 'itu_glass' fired everywhere",
     "Added _tag() helper converting 'nan'/'none' strings to empty string"),
]

for i, (fail, cause, fix) in enumerate(failures):
    row, col = divmod(i, 2)
    x = Inches(0.25 + col * 6.55)
    y = Inches(1.0  + row * 2.1)
    rect(s, x, y, Inches(6.3), Inches(1.95), C_CARD)
    txt(s, "✗ " + fail, x+Inches(0.12), y+Inches(0.06),
        Inches(6.1), Inches(0.38), size=13, bold=True, color=C_WARN)
    txt(s, "Cause: " + cause, x+Inches(0.12), y+Inches(0.46),
        Inches(6.1), Inches(0.5), size=12, color=C_LIGHT)
    txt(s, "✓ Fix: " + fix, x+Inches(0.12), y+Inches(0.98),
        Inches(6.1), Inches(0.5), size=12, color=C_OK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Coverage & Ray Settings
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, 0, W, Inches(0.9), C_CARD)
heading(s, "Coverage Analysis & Ray Tracing Parameters")

# Left: coverage explanation
bullet_box(s, [
    "Scene: 10.76 × 11.25 km = 121 km²",
    "TX position: (4711, 3074) m from centre",
    "TX near NE edge of scene bbox",
    "",
    "3.6 GHz urban macro cell radius ≈ 2 km",
    "π × 2² / 121 km² ≈ 10.4%",
    "Simulated coverage: 11.0%  ✓ matches",
    "",
    "NOT a bug — single cell on large scene",
    "Drive-test points ARE within covered area",
], Inches(0.25), Inches(1.0), Inches(4.2), Inches(5.8),
   title="Why 11% Coverage?", title_color=C_ACCENT)

# Middle: ray settings
bullet_box(s, [
    "MAX_DEPTH = 10 bounces",
    "  (was 6 → added NLOS at 500–1200 m)",
    "NUM_SAMPLES_CM = 200M",
    "  GPU: 100 runs × 2M rays",
    "NUM_SAMPLES_PS = 10M / batch",
    "  Path solver per receiver batch",
    "scat_keep_prob = 0.001",
    "  Scatter rays: 1 in 1000 kept",
    "GRID_SIZE_M = 10 m",
    "  167 rays/cell (borderline for NLOS)",
], Inches(4.6), Inches(1.0), Inches(4.2), Inches(5.8),
   title="Ray Tracing Settings", title_color=C_WARN)

# Right: scatter finding
bullet_box(s, [
    "With scatter:    mean=−55.6 dBm",
    "Without scatter: mean=−55.5 dBm",
    "Scatter impact std = 7.06 dB",
    "  (±7 dB per-cell redistribution)",
    "mean ≈ 0 dB = energy conservation ✓",
    "",
    "Scatter active only after Cell 33",
    "  (re-apply S coeff after scene load)",
    "CAL_MAX_DEPTH must match MAX_DEPTH",
    "  (was 6, should be 10)",
], Inches(8.9), Inches(1.0), Inches(4.1), Inches(5.8),
   title="Scatter Impact", title_color=C_OK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Scene Enhancements Summary
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, 0, W, Inches(0.9), C_CARD)
heading(s, "Scene Enhancements — Final Architecture")

enhancements = [
    ("AWS Elevation Tiles", C_ACCENT,
     ["Zoom-14 GeoTIFF tiles (~5.8 m/px)", "72 tiles downloaded & mosaiced",
      "Bilinear interpolation for smooth mesh", "DEM range: −34 to +102 m local",
      "500×500 terrain grid (~9 MB PLY)"]),
    ("OSM + Overture Buildings", C_WARN,
     ["61,985 OSM building footprints", "70,765 Overture Maps entries",
      "22,558 real heights enriched", "Fallback: levels×3.5 m or 8 m default",
      "47,592 buildings processed → 8 PLY files"]),
    ("Meta AI Canopy Height", C_OK,
     ["Quadkey zoom-9 Web Mercator tiles", "1 m resolution CHM (windowed read)",
      "Height range: 1–53 m in Nottingham", "115,497 trees at 20 m grid spacing",
      "8-sided cylinders, binary PLY (87 MB)"]),
    ("ITU-R P.2040-2 Materials", C_ACCENT,
     ["itu_brick S=0.30, itu_concrete S=0.40", "itu_glass S=0.08, itu_metal S=0.05",
      "itu_vegetation S=0.75 (highest scatter)", "itu_urban_ground ε=5 σ=0.5",
      "LambertianPattern on all materials"]),
]

for i, (title, color, items) in enumerate(enhancements):
    x = Inches(0.25 + (i % 2) * 6.55)
    y = Inches(1.0  + (i // 2) * 3.1)
    rect(s, x, y, Inches(6.3), Inches(2.9), C_CARD)
    rect(s, x, y, Inches(6.3), Inches(0.5), color)
    txt(s, title, x+Inches(0.15), y+Inches(0.06), Inches(6.0), Inches(0.4),
        size=16, bold=True, color=C_BG)
    cy = y + Inches(0.58)
    for item in items:
        txt(s, "• " + item, x+Inches(0.15), cy, Inches(6.0), Inches(0.35),
            size=13, color=C_LIGHT)
        cy += Inches(0.38)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Issues & Remaining Steps
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, C_BG)
rect(s, 0, 0, W, Inches(0.9), C_CARD)
heading(s, "Remaining Issues & Next Steps")

bullet_box(s, [
    "✗ CAL_MAX_DEPTH=6 ≠ MAX_DEPTH=10  →  calibrated S values are for depth=6",
    "✗ TREE_GRID_SPACING was 5 m → 1.8M trees → 1.5 GB PLY → kernel OOM",
    "✗ FLAT_TERRAIN=True still in Cell 0c  →  TX at z=17m not z=25m",
    "✗ Scatter S=0 if Cell 33 not re-run after kernel restart",
    "✗ Scene centred on bbox not TX  →  TX at (4711, 3074) near NE edge",
], Inches(0.25), Inches(1.0), Inches(12.8), Inches(2.4),
   title="Known Issues", title_color=C_WARN, item_size=14)

bullet_box(s, [
    "▶ Set FLAT_TERRAIN=False in Cell 0c → TX Z should show ≈25 m",
    "▶ Run Cell 33 after every scene load → confirm S≠0 in Cell 34",
    "▶ Fix CAL_MAX_DEPTH=10 → re-run calibration sweep for valid S values",
    "▶ Run Cell 9b → get bias/RMSE vs 4001 Ofcom drive-test points",
    "▶ Consider re-centering scene bbox on TX location for better coverage %",
    "▶ Increase scat_keep_prob 0.001→0.01 for better NLOS scatter at 500m+",
], Inches(0.25), Inches(3.55), Inches(12.8), Inches(2.8),
   title="Next Steps", title_color=C_OK, item_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = '/home/user/Ray-Tracing---Sionna-RT/sionna_rt_project_summary.pptx'
prs.save(out)
print(f'Saved: {out}')
